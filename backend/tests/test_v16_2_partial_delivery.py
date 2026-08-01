from pathlib import Path

from pptx import Presentation

from app.engines.job_engine import run_local_job


def _sample_pptx(path: Path) -> Path:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Architecture test"
    slide.placeholders[1].text = "Conversion must remain isolated."
    prs.save(path)
    return path


def test_pptx_headless_formats_complete(tmp_path: Path):
    source = _sample_pptx(tmp_path / "input.pptx")
    result = run_local_job(
        {"order_number": "TEST-V162-OK", "services": ["conversion"], "conversion": {"formats": ["original", "pdf", "xlsx"]}},
        [(source.name, str(source))],
        tmp_path / "output",
    )
    assert result["state"] == "completed"
    assert result["failure_count"] == 0
    assert {Path(item["path"]).suffix for item in result["outputs"]} == {".pptx", ".pdf", ".xlsx"}


def test_pptx_can_export_structured_csv(tmp_path: Path):
    source = _sample_pptx(tmp_path / "input.pptx")
    result = run_local_job(
        {"order_number": "TEST-V162-PARTIAL", "services": ["conversion"], "conversion": {"formats": ["original", "csv"]}},
        [(source.name, str(source))],
        tmp_path / "output",
    )
    assert result["state"] == "completed"
    assert result["partial_success"] is False
    assert result["failure_count"] == 0
    names = {Path(item["path"]).name for item in result["outputs"]}
    assert "input.pptx" in names
    assert "input.csv" in names
