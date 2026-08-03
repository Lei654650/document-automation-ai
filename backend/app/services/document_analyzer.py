from __future__ import annotations

import csv
import re
import time
import xml.etree.ElementTree as ET
import zipfile
from pathlib import Path
from typing import Any

FORMAT_LABELS = {
    '.pdf': 'PDF', '.xlsx': 'Excel', '.xls': 'Excel', '.docx': 'Word', '.doc': 'Word',
    '.pptx': 'PowerPoint', '.ppt': 'PowerPoint', '.csv': 'CSV', '.png': '图片',
    '.jpg': '图片', '.jpeg': '图片', '.bmp': '图片', '.tif': '图片', '.tiff': '图片',
    '.txt': 'TXT', '.zip': 'ZIP',
}


LANGUAGE_LABELS = {
    'zh': '中文（简体）',
    'zh-TW': '中文（繁体）',
    'vi': '越南语',
    'en': '英语',
    'unknown': '未知',
}

_TRADITIONAL_HINTS = set('體臺灣國學醫療證號處護歷責務廠區聯絡資資料書員門間說應與為於實業產點權機關檔圖個進這種')
_SIMPLIFIED_HINTS = set('体台国学医疗证号处护历责任务厂区联络资资料书员门间说应与为于实业产点权机关档图个进这种')
_VIETNAMESE_CHARS_RE = re.compile(r'[ăâđêôơưĂÂĐÊÔƠƯàáảãạằắẳẵặầấẩẫậèéẻẽẹềếểễệìíỉĩịòóỏõọồốổỗộờớởỡợùúủũụừứửữựỳýỷỹỵ]', re.I)
_VIETNAMESE_WORD_RE = re.compile(r'\b(?:và|của|người|công\s+ty|ngày|tháng|năm|địa\s+chỉ|sức\s+khỏe|giấy|cam\s+kết|quy\+?định|nhân\s+viên|tại|không|được|phải|họ\s+và\s+tên)\b', re.I)
_ENGLISH_WORD_RE = re.compile(r'\b(?:the|and|of|to|for|with|company|employee|agreement|document|information|factory|health|passport|address|date|signature|shall|this|that|from|during|business)\b', re.I)


def _detect_sensitive_categories(text: str, filename: str = '') -> list[str]:
    value = f'{filename}\n{text or ""}'
    patterns = [
        ('身份与护照信息', r'passport|hộ\s*chiếu|护照|身份[证證]|cccd|cmnd'),
        ('健康与医疗信息', r'sức\s*khỏe|health\s*(?:check|certificate)|健康|诊断|chẩn\s*đoán|medical'),
        ('无犯罪与司法记录', r'无犯罪|無犯罪|criminal\s*record|lý\s*lịch\s*tư\s*pháp|xác\s*nhận\s*dân\s*sự'),
        ('地址与联系方式', r'địa\s*chỉ|address|地址|điện\s*thoại|telephone|phone|聯絡電話|联系电话'),
        ('签名、指纹或印章', r'ký\s*tên|signature|签名|簽名|指纹|指模|vân\s*tay|印章|đóng\s*dấu'),
        ('就业与背景材料', r'工作经历|工作經歷|quá\s*trình\s*công\s*tác|employment|xin\s*việc|nhân\s*viên'),
    ]
    return [label for label, pattern in patterns if re.search(pattern, value, re.I)]


def _base(name: str, path: Path) -> dict[str, Any]:
    suffix = path.suffix.lower()
    return {
        'name': name,
        'format': FORMAT_LABELS.get(suffix, suffix.lstrip('.').upper() or '未知'),
        'extension': suffix,
        'size_bytes': path.stat().st_size,
        'details': {},
        'warnings': [],
        'capabilities': [],
        '_text_sample': '',
    }


def _detect_language(text: str) -> dict[str, Any]:
    """Detect every meaningful language in a sample instead of forcing one label.

    The former detector treated any Vietnamese diacritic as decisive and therefore
    collapsed mixed Chinese/Vietnamese/English packets into one language.  This
    scorer keeps per-language evidence and returns a backward-compatible primary
    language together with the complete detected set.
    """
    sample = (text or '').strip()[:120000]
    if not sample:
        return {'code': 'unknown', 'name': '未知', 'confidence': 0.0, 'languages': [], 'scores': {}}

    cjk_chars = re.findall(r'[\u3400-\u9fff]', sample)
    cjk = len(cjk_chars)
    traditional = sum(ch in _TRADITIONAL_HINTS for ch in cjk_chars)
    simplified = sum(ch in _SIMPLIFIED_HINTS for ch in cjk_chars)
    vi_chars = len(_VIETNAMESE_CHARS_RE.findall(sample))
    vi_words = len(_VIETNAMESE_WORD_RE.findall(sample))
    en_words = len(_ENGLISH_WORD_RE.findall(sample))
    latin_words = len(re.findall(r'\b[A-Za-z]{3,}\b', sample))

    raw_scores: dict[str, float] = {}
    if cjk >= 4:
        if traditional >= 2 and traditional >= simplified:
            raw_scores['zh-TW'] = cjk + traditional * 4
            if simplified >= 2:
                raw_scores['zh'] = simplified * 4
        elif simplified >= 2:
            raw_scores['zh'] = cjk + simplified * 4
            if traditional >= 2:
                raw_scores['zh-TW'] = traditional * 4
        else:
            raw_scores['zh'] = float(cjk)
    vi_score = vi_chars * 2.4 + vi_words * 10
    if vi_chars >= 2 or vi_words >= 2:
        raw_scores['vi'] = vi_score
    # Require recognisable English function/domain words so Vietnamese Latin text
    # is not also labelled English merely because it contains unaccented words.
    en_score = en_words * 10 + max(0, latin_words - vi_words * 2) * 0.25
    if en_words >= 3 or (en_words >= 1 and latin_words >= 18):
        raw_scores['en'] = en_score

    if not raw_scores:
        return {'code': 'unknown', 'name': '未知', 'confidence': 0.2, 'languages': [], 'scores': {}}
    maximum = max(raw_scores.values())
    threshold = max(4.0, maximum * 0.12)
    ordered = [code for code, score in sorted(raw_scores.items(), key=lambda item: item[1], reverse=True) if score >= threshold]
    primary = ordered[0]
    confidence = round(min(0.99, raw_scores[primary] / max(1.0, sum(raw_scores.values())) + 0.35), 2)
    return {
        'code': primary,
        'name': LANGUAGE_LABELS[primary],
        'confidence': confidence,
        'languages': [{'code': code, 'name': LANGUAGE_LABELS[code], 'confidence': round(min(0.99, raw_scores[code] / max(1.0, maximum)), 2)} for code in ordered],
        'scores': {code: round(score, 2) for code, score in raw_scores.items()},
    }


def _analyze_pdf(item: dict, path: Path) -> None:
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    pages = len(reader.pages)
    chunks: list[str] = []
    page_details: list[dict[str, Any]] = []
    image_count = 0
    scanned_pages = 0
    text_pages = 0
    table_page_estimate = 0
    page_languages: dict[str, int] = {}

    # Inspect every page for normal customer packets. Very large PDFs are sampled
    # evenly so order creation remains bounded while still seeing the whole file.
    if pages <= 80:
        page_indices = list(range(pages))
    else:
        page_indices = sorted(set([*range(20), *range(max(20, pages - 20), pages), *[round(i * (pages - 1) / 39) for i in range(40)]]))

    for index, page in enumerate(reader.pages):
        text = ''
        if index in page_indices:
            try:
                text = page.extract_text() or ''
            except Exception:
                text = ''
            if text.strip():
                chunks.append(text[:8000])
        page_image_count = 0
        try:
            resources = page.get('/Resources') or {}
            x_objects = resources.get('/XObject') or {}
            x_objects = x_objects.get_object() if hasattr(x_objects, 'get_object') else x_objects
            for value in x_objects.values():
                obj = value.get_object() if hasattr(value, 'get_object') else value
                if obj.get('/Subtype') == '/Image':
                    page_image_count += 1
        except Exception:
            page_image_count = 0
        image_count += page_image_count
        chars = len(text.strip())
        likely_scanned_page = chars < 35 and page_image_count > 0 or chars < 8
        if likely_scanned_page:
            scanned_pages += 1
        else:
            text_pages += 1
        lines = [line for line in text.splitlines() if line.strip()]
        looks_tabular = sum(bool(re.search(r'\S\s{2,}\S|\t', line)) for line in lines) >= 3
        if looks_tabular:
            table_page_estimate += 1
        language = _detect_language(text)
        for entry in language.get('languages', []):
            page_languages[entry['name']] = page_languages.get(entry['name'], 0) + 1
        if index in page_indices:
            page_details.append({
                'page': index + 1,
                'text_chars': chars,
                'image_count': page_image_count,
                'likely_scanned': likely_scanned_page,
                'languages': [entry['name'] for entry in language.get('languages', [])],
                'table_like': looks_tabular,
            })

    sample = '\n'.join(chunks)
    item['_text_sample'] = sample
    sensitive = _detect_sensitive_categories(sample, item.get('name', ''))
    item['details'].update({
        'pages': pages,
        'extractable_text_chars_sample': len(sample),
        'likely_scanned': scanned_pages > 0,
        'scanned_page_count': scanned_pages,
        'text_page_count': text_pages,
        'mixed_scan_text': scanned_pages > 0 and text_pages > 0,
        'encrypted': bool(reader.is_encrypted),
        'image_count': image_count,
        'table_page_count_estimate': table_page_estimate,
        'page_language_summary': page_languages,
        'page_analysis_complete': len(page_indices) == pages,
        'page_details': page_details[:120],
        'contains_sensitive_data': bool(sensitive),
        'sensitive_categories': sensitive,
    })
    item['capabilities'] += ['文本提取', '页数识别', '页面级语言识别', '扫描页识别', '敏感信息类别识别']
    if scanned_pages:
        item['warnings'].append(f'检测到 {scanned_pages} 个扫描或图片页面，建议启用 OCR。')
    if scanned_pages and text_pages:
        item['warnings'].append('该 PDF 同时包含扫描页与可提取文本页，需要混合处理。')
    if sensitive:
        item['warnings'].append('检测到身份、健康、联系信息或签章类敏感材料；处理与日志应采用脱敏策略。')


def _analyze_xlsx(item: dict, path: Path) -> None:
    from openpyxl import load_workbook
    # Large or complex workbooks can take minutes and consume hundreds of MB when
    # loaded in normal mode. During order creation we only need a safe structural
    # preview, so files >= 8 MB are opened in read-only streaming mode.
    streaming = path.stat().st_size >= 8 * 1024 * 1024
    wb = load_workbook(path, read_only=streaming, data_only=False, keep_links=False)
    sheets, samples = [], []
    total_rows = total_cells = formula_count = merged_count = chart_count = image_count = 0
    max_columns = 0
    for ws in wb.worksheets:
        rows, cols = ws.max_row or 0, ws.max_column or 0
        total_rows += rows
        max_columns = max(max_columns, cols)
        merged_count += 0 if streaming else len(ws.merged_cells.ranges)
        chart_count += 0 if streaming else len(getattr(ws, '_charts', []))
        image_count += 0 if streaming else len(getattr(ws, '_images', []))
        non_empty = 0
        for row in ws.iter_rows(min_row=1, max_row=min(rows, 80), min_col=1, max_col=min(cols, 40)):
            for cell in row:
                value = cell.value
                if value not in (None, ''):
                    non_empty += 1
                    if isinstance(value, str):
                        samples.append(value)
                    if isinstance(value, str) and value.startswith('='):
                        formula_count += 1
        total_cells += non_empty
        sheets.append({'name': ws.title, 'rows': rows, 'columns': cols, 'sample_non_empty_cells': non_empty})
    item['_text_sample'] = '\n'.join(samples[:500])
    item['details'].update({
        'sheet_count': len(sheets), 'sheets': sheets[:30], 'total_rows': total_rows,
        'max_columns': max_columns, 'sample_non_empty_cells': total_cells,
        'formula_count_sample': formula_count, 'merged_range_count': merged_count,
        'chart_count': chart_count, 'image_count': image_count,
        'analysis_mode': 'streaming_preview' if streaming else 'full_structure',
    })
    item['capabilities'] += ['工作表识别', '公式识别', '合并单元格识别', '图表与图片统计']
    wb.close()


def _analyze_docx(item: dict, path: Path) -> None:
    from docx import Document
    doc = Document(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    table_cells = [cell.text for table in doc.tables for row in table.rows for cell in row.cells if cell.text.strip()]
    sections = len(doc.sections)
    item['_text_sample'] = '\n'.join((paragraphs + table_cells)[:800])
    page_count = 0
    try:
        with zipfile.ZipFile(path) as package:
            if 'docProps/app.xml' in package.namelist():
                root = ET.fromstring(package.read('docProps/app.xml'))
                page_count = int(root.findtext('{http://schemas.openxmlformats.org/officeDocument/2006/extended-properties}Pages') or 0)
    except (ValueError, ET.ParseError, zipfile.BadZipFile):
        page_count = 0
    item['details'].update({
        'paragraph_count': len(doc.paragraphs), 'non_empty_paragraph_count': len(paragraphs),
        'table_count': len(doc.tables), 'image_count': len(doc.inline_shapes), 'section_count': sections,
        'heading_count': sum(1 for p in doc.paragraphs if p.style and p.style.name.lower().startswith('heading')),
        'page_count': page_count,
    })
    item['capabilities'] += ['段落识别', '表格识别', '图片统计', '标题结构识别']


def _analyze_pptx(item: dict, path: Path) -> None:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(path))
    slides = []
    all_text = []
    totals = {'text_shape_count': 0, 'table_count': 0, 'picture_count': 0, 'chart_count': 0, 'group_shape_count': 0}
    for index, slide in enumerate(prs.slides, start=1):
        info = {'slide': index, 'text_shapes': 0, 'tables': 0, 'pictures': 0, 'charts': 0, 'groups': 0, 'text_chars': 0}
        for shape in slide.shapes:
            if getattr(shape, 'has_text_frame', False):
                text = shape.text or ''
                if text.strip():
                    all_text.append(text)
                    info['text_shapes'] += 1
                    info['text_chars'] += len(text)
            if getattr(shape, 'has_table', False):
                info['tables'] += 1
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            all_text.append(cell.text)
            if getattr(shape, 'has_chart', False):
                info['charts'] += 1
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                info['pictures'] += 1
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                info['groups'] += 1
        for key, total_key in [('text_shapes','text_shape_count'),('tables','table_count'),('pictures','picture_count'),('charts','chart_count'),('groups','group_shape_count')]:
            totals[total_key] += info[key]
        slides.append(info)
    item['_text_sample'] = '\n'.join(all_text[:1200])
    item['details'].update({
        'slide_count': len(prs.slides),
        **totals,
        'slide_width_inches': round(prs.slide_width / 914400, 2),
        'slide_height_inches': round(prs.slide_height / 914400, 2),
        'slides': slides[:100],
        'extractable_text_chars': len(item['_text_sample']),
    })
    item['capabilities'] += ['幻灯片识别', '文本框提取', '表格提取', '图片统计', '图表统计', '分组图形统计']
    if totals['group_shape_count']:
        item['warnings'].append('检测到分组图形；复杂 SmartArt/组合对象在后续翻译导出时需要版式复核。')


def _analyze_csv(item: dict, path: Path) -> None:
    rows = 0; max_cols = 0; samples = []
    with path.open('r', encoding='utf-8-sig', errors='replace', newline='') as f:
        for row in csv.reader(f):
            rows += 1; max_cols = max(max_cols, len(row)); samples.extend(row[:20])
    item['_text_sample'] = '\n'.join(samples[:1000])
    item['details'].update({'rows': rows, 'columns': max_cols})
    item['capabilities'] += ['行列结构识别', '文本提取']


def _analyze_txt(item: dict, path: Path) -> None:
    text = path.read_text(encoding='utf-8-sig', errors='replace')
    item['_text_sample'] = text[:30000]
    item['details'].update({
        'character_count': len(text),
        'line_count': len(text.splitlines()),
    })
    item['capabilities'] += ['纯文本提取', '语言识别']


def _analyze_image(item: dict, path: Path) -> None:
    from PIL import Image
    with Image.open(path) as im:
        item['details'].update({'width': im.width, 'height': im.height, 'mode': im.mode, 'image_format': im.format})
    item['warnings'].append('图片文字处理通常需要 OCR。')
    item['capabilities'] += ['图像元数据识别', '可进入 OCR 流程']


def _analyze_zip(item: dict, path: Path) -> None:
    with zipfile.ZipFile(path) as z:
        names = [n for n in z.namelist() if not n.endswith('/')]
    extensions: dict[str, int] = {}
    for name in names:
        ext = Path(name).suffix.lower() or '(无扩展名)'
        extensions[ext] = extensions.get(ext, 0) + 1
    item['details'].update({'contained_file_count': len(names), 'extension_summary': extensions, 'sample_files': names[:50]})
    item['capabilities'] += ['压缩包目录识别', '内部格式统计']


def _category(files: list[dict], requirements: str) -> str:
    text = (' '.join(x['name'] for x in files) + ' ' + requirements + ' ' + ' '.join(x.get('_text_sample','')[:8000] for x in files)).lower()
    if any(key in text for key in ['passport', 'hộ chiếu', '护照', '健康', 'sức khỏe', '工作经历', 'quá trình công tác', '保密承诺', 'cam kết bảo mật']):
        return '人员证件与合规材料'
    rules = [
        (['invoice','发票'], '发票'), (['quotation','quote','报价'], '报价单'), (['contract','合同'], '合同'),
        (['bom','物料清单'], 'BOM'), (['sop','作业指导'], 'SOP'), (['plc','i/o','io list'], 'PLC I/O'),
        (['manual','说明书'], '说明书'), (['assy line','工站','风险','对策'], '产线问题与风险报告'),
    ]
    for keys, label in rules:
        if any(k in text for k in keys): return label
    formats = {x['format'] for x in files}
    if 'PowerPoint' in formats: return '演示文稿'
    if formats <= {'Excel','CSV'}: return '结构化表格数据'
    if '图片' in formats: return '图片/扫描资料'
    return '通用文档'


def _industry(files: list[dict], requirements: str) -> str:
    text = (' '.join(x['name'] for x in files) + ' ' + requirements + ' ' + ' '.join(x.get('_text_sample', '')[:10000] for x in files)).lower()
    if any(key in text for key in ['passport', 'hộ chiếu', '护照', '工作经历', 'quá trình công tác', '员工', 'nhân viên', '驻厂', '駐廠']):
        return '人事与合规'
    rules = [
        (['patient', 'clinical', 'diagnosis', 'medical', 'hospital', 'sức khỏe', '医疗', '患者', '临床', '医院', '健康'], '医疗'),
        (['vehicle', 'automotive', 'engine', 'ecu', 'can bus', '汽车', '车辆', '发动机', '整车'], '汽车'),
        (['contract', 'agreement', 'legal', 'liability', '合同', '协议', '条款', '法律'], '法律'),
        (['invoice', 'balance sheet', 'financial', 'accounting', '发票', '财务', '会计', '审计'], '金融与财务'),
        (['training', 'course', 'curriculum', 'education', '培训', '课程', '教材', '教育'], '教育'),
        (['software', 'database', 'server', 'api', 'cloud', '代码', '软件', '数据库', '服务器'], 'IT'),
        (['solar', 'wind power', 'battery', 'energy', 'oil', 'gas', '电力', '能源', '光伏', '风电', '电池'], '能源'),
        (['construction', 'building', 'architecture', 'bim', '建筑', '施工', '工程图', '土木'], '建筑'),
        (['logistics', 'warehouse', 'shipping', 'freight', 'inventory', '物流', '仓库', '运输', '库存'], '物流'),
        (['pcb', 'semiconductor', 'electronics', 'chip', 'circuit', '电子', '芯片', '电路', '元器件'], '电子'),
        (['plc', 'hmi', 'scada', 'servo', '变频器', '自动化', 'i/o'], '工业自动化'),
        (['bom', 'assembly', 'process sheet', '工艺', '制造', '产线', '零件', '物料'], '制造业'),
        (['marketing', 'campaign', '品牌推广', '营销'], '市场营销'),
    ]
    for keys, label in rules:
        if any(key in text for key in keys):
            return label
    return '通用'


def analyze_order_files(paths: list[tuple[str,str]], services: list[str], requirements: str, translation: dict) -> dict:
    started_at = time.perf_counter()
    files = []
    for name, raw in paths:
        path = Path(raw); item = _base(name, path)
        try:
            ext = path.suffix.lower()
            if ext == '.pdf': _analyze_pdf(item, path)
            elif ext == '.xlsx': _analyze_xlsx(item, path)
            elif ext == '.docx': _analyze_docx(item, path)
            elif ext == '.pptx': _analyze_pptx(item, path)
            elif ext == '.csv': _analyze_csv(item, path)
            elif ext == '.txt': _analyze_txt(item, path)
            elif ext in {'.png','.jpg','.jpeg','.bmp','.tif','.tiff'}: _analyze_image(item, path)
            elif ext == '.zip': _analyze_zip(item, path)
            elif ext in {'.xls','.doc','.ppt'}:
                item['warnings'].append('旧版 Office 二进制格式仅做基础识别；建议转换为 XLSX、DOCX 或 PPTX 后进行深度分析。')
        except Exception as exc:
            item['warnings'].append(f'深度分析未完成：{type(exc).__name__}: {str(exc)[:120]}')
        if 'ocr' in services and item['format'] == '图片':
            item['warnings'] = [w for w in item['warnings'] if '需要 OCR' not in w]
            item['details']['ocr_status'] = '已启用'
        language = _detect_language(item.get('_text_sample',''))
        item['details']['detected_language'] = language
        if not item['details'].get('sensitive_categories'):
            sensitive = _detect_sensitive_categories(item.get('_text_sample', ''), item.get('name', ''))
            item['details']['contains_sensitive_data'] = bool(sensitive)
            item['details']['sensitive_categories'] = sensitive
        files.append(item)

    workflow = ['接收并校验文件', '识别文件格式与内部结构', '检测文档主要语言与行业']
    formats = {x['format'] for x in files}
    if 'PowerPoint' in formats: workflow.append('PowerPoint 对象解析（幻灯片、文字、表格、图片、图表）')
    if 'ocr' in services or any(x['details'].get('likely_scanned') or x['format']=='图片' for x in files): workflow.append('OCR 文字与表格识别')
    if 'data_cleanup' in services: workflow.append('数据清理与结构化')
    if 'translation' in services:
        target = translation.get('target_language','目标语言')
        workflow.append(f'文档翻译（目标：{target}）')
    if 'layout_preserve' in services: workflow.append('版式还原与排版优化')
    outputs = [x.replace('output_','') for x in services if x.startswith('output_')]
    if outputs: workflow.append('生成输出文件：' + ', '.join(outputs))
    if 'manual_review' in services: workflow.append('人工质量复核')
    workflow.append('交付文件')

    total = sum(x['size_bytes'] for x in files)
    total_pages = 0
    total_images = 0
    total_tables = 0
    page_count_complete = True
    image_count_complete = True
    table_count_complete = True
    ocr_required = False
    ocr_languages: set[str] = set()
    for item in files:
        details = item.get('details', {})
        fmt = item.get('format')
        page_value = details.get('pages')
        if page_value is None and fmt == 'PowerPoint':
            page_value = details.get('slide_count')
        if page_value is None and fmt == 'Word':
            page_value = details.get('page_count')
        if page_value is None and fmt == '图片':
            page_value = 1
        if page_value:
            total_pages += int(page_value)
        elif fmt in {'PDF', 'Word', 'PowerPoint', '图片'}:
            page_count_complete = False

        image_value = details.get('image_count')
        if image_value is None and fmt == 'PowerPoint':
            image_value = details.get('picture_count')
        if image_value is not None:
            total_images += int(image_value or 0)
        elif fmt in {'PDF', 'Word', 'PowerPoint'}:
            image_count_complete = False

        table_value = details.get('table_count')
        if table_value is None and fmt == 'PDF':
            table_value = details.get('table_page_count_estimate')
        if table_value is not None:
            total_tables += int(table_value or 0)
        elif fmt in {'PDF', '图片'}:
            table_count_complete = False

        needs_ocr = bool(details.get('likely_scanned')) or fmt == '图片'
        if needs_ocr:
            ocr_required = True
            language_info = details.get('detected_language', {})
            language_entries = language_info.get('languages') or []
            if language_entries:
                for entry in language_entries:
                    if entry.get('name') and entry.get('name') != '未知':
                        ocr_languages.add(entry['name'])
            else:
                language_name = language_info.get('name')
                if language_name and language_name != '未知':
                    ocr_languages.add(language_name)
    warnings = [w for x in files for w in x['warnings']]
    detected_language_names: list[str] = []
    language_codes: list[str] = []
    sensitive_categories: list[str] = []
    scanned_pages = text_pages = 0
    for item in files:
        language = item.get('details', {}).get('detected_language', {})
        entries = language.get('languages') or ([{'code': language.get('code'), 'name': language.get('name')}] if language.get('name') not in {None, '未知'} else [])
        for entry in entries:
            name = entry.get('name')
            code = entry.get('code')
            if name and name != '未知' and name not in detected_language_names:
                detected_language_names.append(name)
            if code and code != 'unknown' and code not in language_codes:
                language_codes.append(code)
        details = item.get('details', {})
        scanned_pages += int(details.get('scanned_page_count') or (1 if details.get('likely_scanned') else 0))
        text_pages += int(details.get('text_page_count') or 0)
        for category_name in details.get('sensitive_categories') or []:
            if category_name not in sensitive_categories:
                sensitive_categories.append(category_name)

    languages = detected_language_names or ['未知']
    mixed_language = len(detected_language_names) > 1
    object_count = sum(
        (x['details'].get('text_shape_count') or 0)
        + (x['details'].get('table_count') or x['details'].get('table_page_count_estimate') or 0)
        + (x['details'].get('picture_count') or x['details'].get('image_count') or 0)
        for x in files
    )
    complexity_score = 0
    if len(files) > 3 or total > 20 * 1024 * 1024 or object_count > 100 or total_pages > 5:
        complexity_score += 1
    if len(files) > 10 or total > 100 * 1024 * 1024 or object_count > 500 or total_pages > 40:
        complexity_score += 1
    if mixed_language:
        complexity_score += 1
    if scanned_pages and text_pages:
        complexity_score += 1
    if len(sensitive_categories) >= 2:
        complexity_score += 1
    complexity = '高' if complexity_score >= 3 else '中' if complexity_score >= 1 else '低'
    category = _category(files, requirements)
    industry = _industry(files, requirements)
    for item in files:
        item.pop('_text_sample', None)
    duration_ms = max(1, round((time.perf_counter() - started_at) * 1000))
    return {
        'engine_version': '45.1-multilingual-document-analyzer',
        'status': 'completed',
        'file_count': len(files),
        'total_size_bytes': total,
        'input_formats': sorted(formats),
        'source_extensions': sorted({Path(item.get('name') or '').suffix.lower().lstrip('.') for item in files if Path(item.get('name') or '').suffix}),
        'mixed_input_formats': len(formats) > 1,
        'format_flags': {
            'has_pdf': 'PDF' in formats,
            'has_word': 'Word' in formats,
            'has_spreadsheet': bool(formats & {'Excel', 'CSV'}),
            'has_presentation': 'PowerPoint' in formats,
            'has_images': '图片' in formats,
        },
        'detected_languages': languages,
        'detected_language_codes': language_codes,
        'mixed_language': mixed_language,
        'document_category': category,
        'industry': industry,
        'complexity': complexity,
        'files': files,
        'recommended_workflow': workflow,
        'total_pages': total_pages,
        'page_count_complete': page_count_complete,
        'scanned_page_count': scanned_pages,
        'text_page_count': text_pages,
        'mixed_scan_text': bool(scanned_pages and text_pages),
        'total_images': total_images,
        'image_count_complete': image_count_complete,
        'total_tables': total_tables,
        'table_count_complete': table_count_complete,
        'ocr_required': ocr_required,
        'ocr_languages': sorted(ocr_languages),
        'contains_sensitive_data': bool(sensitive_categories),
        'sensitive_categories': sensitive_categories,
        'privacy_notice': '检测到敏感材料时，界面与日志仅展示类别，不展示证件号码、健康内容或联系方式。' if sensitive_categories else '',
        'analysis_duration_ms': duration_ms,
        'warnings': warnings or ['未发现明显风险。'],
        'summary': f'已深度识别 {len(files)} 个文件；格式：{", ".join(sorted(formats))}；类别：“{category}”；行业：“{industry}”；语言：{", ".join(languages)}；复杂度：{complexity}。'
    }
