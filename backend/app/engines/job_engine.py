from __future__ import annotations

import json
import copy
import re
import shutil
import os
import logging
import time
import tempfile
import zipfile
import xml.etree.ElementTree as ET
from concurrent.futures import ThreadPoolExecutor, as_completed
from threading import Lock
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Callable, Iterable

from docx import Document
from docx.text.paragraph import Paragraph
from docx.oxml.ns import qn
from docx.shared import Pt, Inches, Cm
from openpyxl import load_workbook
from openpyxl.worksheet.pagebreak import Break
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pypdf import PdfReader

try:
    import fitz  # PyMuPDF, used to render scanned PDF pages for OCR
except Exception:  # pragma: no cover
    fitz = None

from .ocr_engine import capability as ocr_capability, extract_text_from_image
from .translation_engine import TranslationClient, capability as translation_capability
from .conversion_engine import convert_outputs

ProgressCallback = Callable[[int, str, str], None]
LOGGER = logging.getLogger("document_automation.jobs")


def now() -> str:
    return datetime.now(timezone.utc).isoformat()


def build_plan(order: dict[str, Any]) -> list[dict[str, Any]]:
    services = order.get("services") or []
    steps: list[dict[str, Any]] = [
        {"id": "validate", "label": "Validate source files", "required": True},
        {"id": "analyze", "label": "Analyze document structure", "required": True},
    ]
    if "ocr" in services:
        steps.append({"id": "ocr", "label": "OCR scanned content", "required": True})
    if "translation" in services:
        steps.append({"id": "translation", "label": "AI automatic translation", "required": True})
    if "data_cleanup" in services:
        steps.append({"id": "cleanup", "label": "智能数据整理", "required": True})
    if "enterprise_analysis" in services:
        steps.append({"id": "analysis_report", "label": "企业数据分析", "required": True})
    if "conversion" in services:
        steps.append({"id": "conversion", "label": "Convert selected output formats", "required": True})
    if "layout_preserve" in services:
        steps.append({"id": "layout", "label": "Preserve or recover layout", "required": True})
    steps.append({"id": "quality", "label": "Validate processing quality", "required": True})
    steps.append({"id": "export", "label": "Generate delivery files", "required": True})
    if "manual_review" in services:
        steps.append({"id": "review", "label": "Manual quality review", "required": True})
    return steps


def _resolve_job_outcome(successful_output_count: int, failure_count: int, manual_review: bool = False) -> tuple[str, str]:
    """Return the single authoritative terminal state and user-facing message."""
    if successful_output_count <= 0:
        if failure_count > 0:
            return "failed", f"处理失败：0 个文件可交付，{failure_count} 项失败"
        return "failed", "处理失败：没有生成可交付文件"
    if failure_count > 0:
        return "partial_completed", f"部分完成：成功交付 {successful_output_count} 个文件，{failure_count} 项失败"
    if manual_review:
        return "quality_review", f"处理完成：{successful_output_count} 个文件等待人工验收"
    return "completed", f"处理完成：成功交付 {successful_output_count} 个文件"


def _update(callback: ProgressCallback | None, progress: int, step: str, message: str) -> None:
    if callback:
        callback(progress, step, message)


def _translate_pptx(source: Path, destination: Path, client: TranslationClient, callback: ProgressCallback | None) -> int:
    presentation = Presentation(source)
    translated = 0
    total = max(1, len(presentation.slides))
    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    original = "".join(run.text for run in paragraph.runs) or paragraph.text
                    if original.strip():
                        value = client.translate(original)
                        if paragraph.runs:
                            paragraph.runs[0].text = value
                            for run in paragraph.runs[1:]:
                                run.text = ""
                        else:
                            paragraph.text = value
                        translated += 1
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            original = "".join(run.text for run in paragraph.runs) or paragraph.text
                            if original.strip():
                                value = client.translate(original)
                                if paragraph.runs:
                                    paragraph.runs[0].text = value
                                    for run in paragraph.runs[1:]:
                                        run.text = ""
                                else:
                                    paragraph.text = value
                                translated += 1
        _update(callback, 35 + int(slide_index / total * 45), "translation", f"Translated PowerPoint slide {slide_index}/{total}")
    presentation.save(destination)
    return translated


TARGET_SUFFIXES = {
    "zh": "zh-CN", "zh_tw": "zh-TW", "vi": "vi", "en": "en",
    "ja": "ja", "ko": "ko", "th": "th", "fr": "fr", "de": "de",
    "es": "es", "pt": "pt", "ru": "ru", "ar": "ar",
}


def _target_suffix(client: TranslationClient | None) -> str:
    if client is None:
        return "processed"
    raw = str(getattr(client, "target_language_code", "") or "").strip().lower()
    if not raw:
        raw = str(getattr(client, "target_language", "") or "").strip().lower()
    return TARGET_SUFFIXES.get(raw, re.sub(r"[^a-z0-9-]+", "-", raw).strip("-") or "translated")


def _paragraph_text(paragraph) -> str:
    """Return all visible text, including text nested in hyperlinks/text boxes."""
    nodes = paragraph._p.xpath('.//w:t')
    xml_text = ''.join(node.text or '' for node in nodes)
    return xml_text or paragraph.text or ''


def _preferred_font_for_target(target_code: str) -> str | None:
    code = (target_code or "").lower()
    if code in {"zh", "zh_tw"}:
        return "Microsoft YaHei"
    if code == "ja":
        return "Yu Gothic"
    if code == "ko":
        return "Malgun Gothic"
    if code in {"ar"}:
        return "Arial"
    return None


def _apply_target_font(run, target_code: str) -> None:
    """Attach a real installed font and remove theme overrides that cause □ glyphs."""
    font_name = _preferred_font_for_target(target_code)
    if not font_name:
        return
    run.font.name = font_name
    r_pr = run._element.get_or_add_rPr()
    r_fonts = r_pr.rFonts
    if r_fonts is None:
        r_fonts = r_pr.get_or_add_rFonts()
    for attr in ("asciiTheme", "hAnsiTheme", "eastAsiaTheme", "cstheme"):
        key = qn(f"w:{attr}")
        if key in r_fonts.attrib:
            del r_fonts.attrib[key]
    for attr in ("eastAsia", "ascii", "hAnsi", "cs"):
        r_fonts.set(qn(f"w:{attr}"), font_name)
    r_pr.set(qn("w:hint"), "eastAsia")


def _configure_document_fonts(document, target_code: str) -> None:
    """Set CJK-compatible defaults at document and style level for Word and WPS."""
    font_name = _preferred_font_for_target(target_code)
    if not font_name:
        return
    for style_name in ("Normal", "Title", "Subtitle", "Heading 1", "Heading 2", "Heading 3"):
        try:
            style = document.styles[style_name]
        except KeyError:
            continue
        style.font.name = font_name
        r_pr = style.element.get_or_add_rPr()
        r_fonts = r_pr.rFonts
        if r_fonts is None:
            r_fonts = r_pr.get_or_add_rFonts()
        for attr in ("asciiTheme", "hAnsiTheme", "eastAsiaTheme", "cstheme"):
            key = qn(f"w:{attr}")
            if key in r_fonts.attrib:
                del r_fonts.attrib[key]
        for attr in ("eastAsia", "ascii", "hAnsi", "cs"):
            r_fonts.set(qn(f"w:{attr}"), font_name)


def _replace_paragraph_text(paragraph, value: str, target_code: str = "") -> None:
    """Replace text without flattening paragraph-level layout or non-text XML.

    Word frequently stores a sentence across several runs and may place text inside
    hyperlinks or text boxes. Updating the underlying ``w:t`` nodes preserves the
    surrounding run formatting, hyperlinks, drawings, bookmarks and field codes.
    """
    text_nodes = list(paragraph._p.xpath('.//w:t'))
    if text_nodes:
        anchor = next((node for node in text_nodes if node.text), text_nodes[0])
        for node in text_nodes:
            node.text = value if node is anchor else ''
        # Apply an East-Asian-capable font to the run containing the anchor node.
        parent = anchor.getparent()
        for run in paragraph.runs:
            if run._element is parent:
                _apply_target_font(run, target_code)
                break
    else:
        run = paragraph.add_run(value)
        _apply_target_font(run, target_code)


def _validate_docx(path: Path) -> dict[str, Any]:
    """Reopen and inspect the generated DOCX before it is delivered."""
    document = Document(path)
    text_blocks = []
    for paragraph in _iter_docx_paragraphs(document):
        text = _paragraph_text(paragraph).strip()
        if text:
            text_blocks.append(text)
    joined = "\n".join(text_blocks)
    bad_chars = {"\ufffd": joined.count("\ufffd"), "\u25a1": joined.count("\u25a1")}
    bad_chars = {key: value for key, value in bad_chars.items() if value}
    if bad_chars:
        raise RuntimeError(f"Generated DOCX contains invalid replacement characters: {bad_chars}")
    return {
        "paragraphs": len(document.paragraphs),
        "tables": len(document.tables),
        "sections": len(document.sections),
        "text_blocks": len(text_blocks),
        "characters": len(joined),
        "invalid_character_count": 0,
    }

def _iter_docx_paragraphs(document):
    """Yield each physical DOCX paragraph exactly once.

    ``id(lxml_element)`` is not a safe de-duplication key because lxml may create
    short-lived Python wrappers for the same XML node and may later reuse their
    object ids.  That caused linked headers/footers to be counted twice in the
    source and then fewer times after reopening the output, producing false
    errors such as ``source had 13 text blocks but output has 11``.

    A stable key is the OOXML part name plus the node's absolute XPath inside
    that part.  It also preserves distinct paragraphs that happen to contain the
    same text.
    """
    seen: set[tuple[str, str]] = set()

    def emit(root, parent, part_name: str):
        tree = root.getroottree()
        for element in root.xpath('.//w:p'):
            key = (part_name, tree.getpath(element))
            if key in seen:
                continue
            seen.add(key)
            yield Paragraph(element, parent)

    body_part = str(getattr(document.part, 'partname', '/word/document.xml'))
    yield from emit(document.element.body, document, body_part)

    for section in document.sections:
        containers = (
            section.header, section.footer,
            section.first_page_header, section.first_page_footer,
            section.even_page_header, section.even_page_footer,
        )
        for container in containers:
            part_name = str(getattr(container.part, 'partname', ''))
            yield from emit(container._element, container, part_name)


def _translate_docx(source: Path, destination: Path, client: TranslationClient, callback: ProgressCallback | None) -> dict[str, Any]:
    document = Document(source)
    _configure_document_fonts(document, client.target_language_code)
    paragraphs = list(_iter_docx_paragraphs(document))
    translatable = [p for p in paragraphs if _paragraph_text(p).strip()]
    total = max(1, len(translatable))
    translated = 0
    skipped = 0
    for index, paragraph in enumerate(translatable, start=1):
        original = _paragraph_text(paragraph)
        value = client.translate(original)
        # A provider must never be allowed to erase a visible source block.
        # Preserve the original and report it as skipped when the response is blank.
        if not str(value or '').strip():
            value = original
        if value != original:
            _replace_paragraph_text(paragraph, value, client.target_language_code)
            translated += 1
        else:
            skipped += 1
        if index % 5 == 0 or index == total:
            _update(callback, 35 + int(index / total * 45), "translation", f"Translated Word content {index}/{total}")

    # Preserve the original document package and add traceable metadata.
    document.core_properties.comments = "Translated by Document Automation AI"
    document.save(destination)
    validation = _validate_docx(destination)
    if validation["text_blocks"] != len(translatable):
        raise RuntimeError(
            f"DOCX validation failed: source had {len(translatable)} unique text blocks but output has "
            f"{validation['text_blocks']}."
        )
    return {
        "translated_items": translated,
        "skipped_items": skipped,
        "source_text_blocks": len(translatable),
        "translation_coverage": round(translated / max(1, len(translatable)) * 100, 2),
        "validation": validation,
    }


def _excel_column_number(ref: str) -> int:
    match = re.match(r"([A-Z]+)", ref or "")
    if not match:
        return 1
    value = 0
    for ch in match.group(1):
        value = value * 26 + ord(ch) - 64
    return value


def _excel_column_name(number: int) -> str:
    number = max(1, int(number))
    result = []
    while number:
        number, remainder = divmod(number - 1, 26)
        result.append(chr(65 + remainder))
    return "".join(reversed(result))


def _prepare_xlsx_for_processing(source: Path, callback: ProgressCallback | None = None, step: str = "translation") -> tuple[Path, Path | None]:
    """Create a compact temporary XLSX when a worksheet contains millions of empty styled cells.

    Some customer workbooks declare A1:XFDxxxx and store one self-closing ``<c/>``
    node for nearly every empty formatted cell.  openpyxl must materialize those
    nodes and can appear frozen for many minutes.  The compact copy removes only
    empty cell nodes from pathological sheets, keeps all cells containing values,
    formulas or inline strings, and updates the worksheet dimension.  The source
    file is never modified.
    """
    threshold = max(20 * 1024 * 1024, int(os.getenv("XLSX_XML_COMPACT_THRESHOLD_MB", "80")) * 1024 * 1024)
    try:
        with zipfile.ZipFile(source, "r") as archive:
            worksheet_infos = [info for info in archive.infolist() if info.filename.startswith("xl/worksheets/") and info.filename.endswith(".xml")]
            pathological = [info for info in worksheet_infos if info.file_size >= threshold]
            if not pathological:
                return source, None
            _update(callback, 5, step, f"检测到超大 Excel 工作表，正在压缩 {len(pathological)} 个异常空白区域")
            temp_handle = tempfile.NamedTemporaryFile(prefix="docai_compact_", suffix=".xlsx", delete=False)
            temp_path = Path(temp_handle.name)
            temp_handle.close()
            with zipfile.ZipFile(temp_path, "w", compression=zipfile.ZIP_DEFLATED, compresslevel=6) as output:
                for info in archive.infolist():
                    data = archive.read(info.filename)
                    if info in pathological:
                        before = len(data)
                        # Remove only truly empty self-closing cells. Cells with
                        # values/formulas/inline strings are never self-closing.
                        data = re.sub(br"<c\b[^>]*/>", b"", data)
                        refs = re.findall(br"<c\b[^>]*\br=\"([A-Z]+[0-9]+)\"[^>]*>", data)
                        if refs:
                            max_row = 1
                            max_col = 1
                            for raw_ref in refs:
                                ref = raw_ref.decode("ascii", "ignore")
                                row_match = re.search(r"(\d+)$", ref)
                                if row_match:
                                    max_row = max(max_row, int(row_match.group(1)))
                                max_col = max(max_col, _excel_column_number(ref))
                            dimension = f'A1:{_excel_column_name(max_col)}{max_row}'.encode("ascii")
                            data = re.sub(br'<dimension\s+ref="[^"]+"\s*/>', b'<dimension ref="' + dimension + b'"/>', data, count=1)
                        LOGGER.info("Compacted worksheet XML: file=%s sheet=%s before=%s after=%s", source.name, info.filename, before, len(data))
                    output.writestr(info, data)
            _update(callback, 10, step, "异常空白区域压缩完成，正在打开实际数据")
            return temp_path, temp_path
    except Exception:
        LOGGER.exception("Failed to compact pathological workbook: %s", source)
        raise




_CJK_RE = re.compile(r"[\u3400-\u9fff]")
_TARGET_WORD_RE = re.compile(r"[A-Za-zÀ-ỹ]+(?:[ '\-][A-Za-zÀ-ỹ]+)*")
_TECH_TOKEN_RE = re.compile(r"^(?:[A-Z]{1,8}\d+[A-Za-z0-9_.:/\-]*|[A-Za-z0-9]+(?:_[A-Za-z0-9]+)+|[A-Z]{2,}[A-Z0-9_.:/\-]*)$")
_TECH_PREFIX_RE = re.compile(
    r"^\s*((?:[A-Z]{1,8}\d+[A-Za-z0-9_.:/\-]*|[A-Za-z0-9]+(?:_[A-Za-z0-9]+)+|[A-Z]{2,}[A-Z0-9_.:/\-]*))"
    r"(?:\s*[|｜]\s*|\s+)(.+)$",
    re.S,
)


def _split_existing_bilingual_text(text: str) -> tuple[str, str] | None:
    """Return the Chinese/source side and existing Latin translation.

    Handles explicit separators, line breaks and historical glued output while
    ignoring PLC/HMI identifiers such as X001, SD0 and Tip0_Title.
    """
    value = str(text or "").strip()
    if not _CJK_RE.search(value):
        return None

    # Prefer explicit separators because they are deterministic.
    for sep in ("——", "—", "\n"):
        if sep in value:
            left, right = value.split(sep, 1)
            if _CJK_RE.search(left) and re.search(r"[A-Za-zÀ-ỹ]", right):
                return left.strip(" \t\r\n:/—-–"), right.strip(" \t\r\n:/—-–")

    last_cjk = max((m.end() for m in _CJK_RE.finditer(value)), default=0)
    for match in _TARGET_WORD_RE.finditer(value, pos=last_cjk):
        token = match.group(0).strip()
        compact = token.replace(" ", "")
        if not token or _TECH_TOKEN_RE.fullmatch(compact):
            continue
        if not (re.search(r"[a-zà-ỹ]", token) or re.search(r"[À-ỹ]", token)):
            continue
        source = value[:match.start()].rstrip(" \t\r\n:/—-–")
        target = value[match.start():].strip(" \t\r\n:/—-–")
        if source and target:
            return source, target
    return None


def _split_technical_prefix(text: str) -> tuple[str, str]:
    """Split a leading PLC/HMI address or identifier from Chinese text."""
    value = str(text or "").strip().replace("｜", "|")
    match = _TECH_PREFIX_RE.match(value)
    if not match:
        return "", value.strip(" |\t\r\n")
    prefix, body = match.group(1).strip(), match.group(2).strip(" |\t\r\n")
    return prefix, body


def _extract_chinese_translation_unit(text: str) -> str:
    existing = _split_existing_bilingual_text(text)
    source_side = existing[0] if existing else str(text or "")
    _, body = _split_technical_prefix(source_side)
    parts = re.findall(r"[\u3400-\u9fff]+", body)
    return " ".join(parts).strip()


def _clean_translation_candidate(source: str, translated: str) -> str:
    candidate = str(translated or "").strip()
    candidate_pair = _split_existing_bilingual_text(candidate)
    if candidate_pair:
        candidate = candidate_pair[1]
    unit = _extract_chinese_translation_unit(source)
    for removable in (str(source or "").strip(), unit):
        if removable:
            candidate = candidate.replace(removable, " ")
    candidate = candidate.replace("|", " ")
    candidate = re.sub(r"\s+", " ", candidate).strip(" \t\r\n:/—-–")
    return candidate


def _normalize_bilingual_value(source: str, translated: str) -> str:
    """Create a readable bilingual cell without leaking internal separators.

    Plain labels use ``中文 —— 越语``. PLC/HMI cells use three lines so the
    address, Chinese description and translation remain visually distinct:
    ``SD0\n前上左安全门\nCửa ...``.
    """
    existing = _split_existing_bilingual_text(source)
    source_side = existing[0] if existing else str(source or "").strip()
    prefix, body = _split_technical_prefix(source_side)
    candidate = _clean_translation_candidate(source, translated)
    if not candidate and existing:
        candidate = existing[1].strip()

    body = body.replace("|", " ").strip()
    if not candidate:
        return "\n".join(part for part in (prefix, body) if part)
    if prefix and _CJK_RE.search(body):
        return f"{prefix}\n{body}\n{candidate}"
    return f"{body} —— {candidate}" if body else candidate


def _apply_excel_multiline_layout(workbook_path: Path) -> None:
    """Enable wrapping and readable row heights for generated multiline cells.

    The update is performed directly in OOXML to avoid openpyxl rewriting the
    workbook. Only cells that contain generated line breaks receive a cloned
    wrapped style; all other workbook structure and styling stays untouched.
    """
    ns = "{http://schemas.openxmlformats.org/spreadsheetml/2006/main}"
    with zipfile.ZipFile(workbook_path, "r") as src:
        payload = {info.filename: src.read(info.filename) for info in src.infolist()}
        infos = list(src.infolist())

    multiline_shared: dict[int, int] = {}
    shared_name = "xl/sharedStrings.xml"
    if shared_name in payload:
        shared_root = ET.fromstring(payload[shared_name])
        for idx, item in enumerate(shared_root.findall(ns + "si")):
            text = "".join(node.text or "" for node in item.iter(ns + "t"))
            if "\n" in text:
                multiline_shared[idx] = max(2, text.count("\n") + 1)
            elif " —— " in text:
                multiline_shared[idx] = 2

    styles_name = "xl/styles.xml"
    if styles_name not in payload:
        return
    styles_root = ET.fromstring(payload[styles_name])
    cell_xfs = styles_root.find(ns + "cellXfs")
    if cell_xfs is None:
        return

    from copy import deepcopy
    wrapped_style_by_key: dict[tuple[int, int], int] = {}

    def wrapped_style(style_id: int, line_count: int) -> int:
        key = (style_id, min(3, line_count))
        if key in wrapped_style_by_key:
            return wrapped_style_by_key[key]
        base = list(cell_xfs)[style_id] if 0 <= style_id < len(cell_xfs) else list(cell_xfs)[0]
        clone = deepcopy(base)
        alignment = clone.find(ns + "alignment")
        if alignment is None:
            alignment = ET.SubElement(clone, ns + "alignment")
        alignment.set("wrapText", "1")
        alignment.set("vertical", "center")
        clone.set("applyAlignment", "1")
        cell_xfs.append(clone)
        new_id = len(cell_xfs) - 1
        wrapped_style_by_key[key] = new_id
        return new_id

    changed_sheets: dict[str, ET.Element] = {}
    for name, data in list(payload.items()):
        if not (name.startswith("xl/worksheets/") and name.endswith(".xml")):
            continue
        root = ET.fromstring(data)
        dirty = False
        for row in root.iter(ns + "row"):
            row_lines = 1
            for cell in row.findall(ns + "c"):
                lines = 1
                if cell.get("t") == "s":
                    value_node = cell.find(ns + "v")
                    if value_node is not None and value_node.text and value_node.text.isdigit():
                        lines = multiline_shared.get(int(value_node.text), 1)
                elif cell.get("t") == "inlineStr":
                    text = "".join(node.text or "" for node in cell.iter(ns + "t"))
                    lines = max(1, text.count("\n") + 1)
                    if lines == 1 and " —— " in text:
                        lines = 2
                if lines > 1:
                    old_style = int(cell.get("s", "0") or 0)
                    cell.set("s", str(wrapped_style(old_style, lines)))
                    row_lines = max(row_lines, lines)
                    dirty = True
            if row_lines > 1:
                current = float(row.get("ht", "0") or 0)
                target = 34.0 if row_lines == 2 else 48.0
                row.set("ht", f"{max(current, target):g}")
                row.set("customHeight", "1")
        if dirty:
            changed_sheets[name] = root

    if not changed_sheets:
        return
    cell_xfs.set("count", str(len(cell_xfs)))
    payload[styles_name] = ET.tostring(styles_root, encoding="utf-8", xml_declaration=True)
    for name, root in changed_sheets.items():
        payload[name] = ET.tostring(root, encoding="utf-8", xml_declaration=True)

    temp = workbook_path.with_suffix(workbook_path.suffix + ".layout.tmp")
    with zipfile.ZipFile(temp, "w", compression=zipfile.ZIP_DEFLATED, compresslevel=6) as out:
        for info in infos:
            out.writestr(info, payload[info.filename])
    temp.replace(workbook_path)


def _translate_reconstructed_xlsx(source: Path, destination: Path, client: TranslationClient, callback: ProgressCallback | None) -> int | None:
    """Fill all dedicated target-language columns in reconstructed workbooks.

    Fixed automation terminology is served from the embedded glossary first;
    only unknown phrases are sent to the configured translation provider.
    """
    try:
        wb = load_workbook(source, data_only=False)
    except Exception:
        return None
    try:
        if '文档概览' not in wb.sheetnames or 'PLC输入信号' not in wb.sheetnames:
            return None
        target_code = getattr(client, 'target_language_code', '')
        if target_code not in {'zh-vi','zh-en'}:
            return None
        target_label = '越南语' if target_code == 'zh-vi' else '英语'
        pairs=[]; sources=[]; seen=set()
        explicit_pairs={
            'PLC输入信号':[(3,4),(5,6),(7,8)],
            '设备清单':[(1,2),(4,5),(6,7),(8,9)],
            '气缸IO配置':[(2,3),(6,7),(8,9),(12,13)],
            '工位结构':[(2,3),(5,6),(7,8)],
            '操作提示':[(2,3),(5,6)],
        }
        for sheet_name,col_pairs in explicit_pairs.items():
            if sheet_name not in wb.sheetnames: continue
            ws=wb[sheet_name]
            for src_col,dst_col in col_pairs:
                if src_col>ws.max_column or dst_col>ws.max_column: continue
                for row in range(2,ws.max_row+1):
                    text=_clean_reconstructed_value(ws.cell(row,src_col).value)
                    if not text or not _CJK_RE.search(text): continue
                    existing=_clean_reconstructed_value(ws.cell(row,dst_col).value)
                    if existing and not _CJK_RE.search(existing): continue
                    pairs.append((ws,row,dst_col,text))
                    if text not in seen: seen.add(text); sources.append(text)
        if not pairs:
            shutil.copy2(source,destination)
            return 0
        _update(callback,24,'translation',f'企业重构表已识别：{len(sources)} 条待补充术语，正在填充独立{target_label}列')
        mapping={}; pending=[]
        for src in sources:
            fixed=_glossary_vi(src) if target_code=='zh-vi' else ''
            if fixed and not _CJK_RE.search(fixed): mapping[src]=fixed
            else: pending.append(src)
        if pending:
            translated=client.translate_many(pending)
            for src,dst in zip(pending,translated):
                value=_clean_translation_candidate(src,str(dst or ''))
                if value and not _CJK_RE.search(value): mapping[src]=value
        failures=[src for src in sources if not mapping.get(src)]
        if failures:
            raise RuntimeError(f'企业重构翻译未完成：仍有 {len(failures)} 条中文没有目标语言，示例：{failures[:5]}')
        for ws,row,col,text in pairs:
            cell=ws.cell(row,col); cell.value=mapping[text]
            cell.font=Font(name='Arial',size=10.5)
            cell.alignment=Alignment(vertical='center',wrap_text=True)
            ws.row_dimensions[row].height=max(ws.row_dimensions[row].height or 26,30)
        destination.parent.mkdir(parents=True,exist_ok=True)
        wb.save(destination)
        _update(callback,78,'translation',f'企业重构翻译完成：{len(pairs)} 个单元格已写入独立{target_label}列')
        return len(pairs)
    finally:
        wb.close()


def _translate_xlsx(source: Path, destination: Path, client: TranslationClient, callback: ProgressCallback | None) -> int:
    """Translate XLSX text in-place inside the OOXML package."""
    reconstructed = _translate_reconstructed_xlsx(source, destination, client, callback)
    if reconstructed is not None:
        return reconstructed
    _update(callback, 10, "translation", "正在读取 Excel 文本索引（不会重排工作表）")
    with zipfile.ZipFile(source, "r") as archive:
        names = set(archive.namelist())
        xml_payloads: dict[str, bytes] = {}
        text_nodes: list[tuple[str, ET.Element, list[ET.Element], str]] = []
        unique_texts: list[str] = []
        seen: set[str] = set()
        ns = "{http://schemas.openxmlformats.org/spreadsheetml/2006/main}"

        def register(package_name: str, root: ET.Element) -> None:
            for parent in root.iter():
                if parent.tag not in {ns + "si", ns + "is"}:
                    continue
                nodes = list(parent.iter(ns + "t"))
                if not nodes:
                    continue
                original = "".join(node.text or "" for node in nodes)
                if not (_should_translate_excel_text(original) or _split_existing_bilingual_text(original)):
                    continue
                text_nodes.append((package_name, parent, nodes, original))
                if original not in seen:
                    seen.add(original)
                    unique_texts.append(original)

        if "xl/sharedStrings.xml" in names:
            data = archive.read("xl/sharedStrings.xml")
            root = ET.fromstring(data)
            xml_payloads["xl/sharedStrings.xml"] = data
            register("xl/sharedStrings.xml", root)
            xml_payloads["xl/sharedStrings.xml:root"] = root  # type: ignore[assignment]

        for name in sorted(n for n in names if n.startswith("xl/worksheets/") and n.endswith(".xml")):
            data = archive.read(name)
            if b"inlineStr" not in data:
                continue
            root = ET.fromstring(data)
            xml_payloads[name] = data
            register(name, root)
            xml_payloads[name + ":root"] = root  # type: ignore[assignment]

        total = len(unique_texts)
        if not total:
            shutil.copy2(source, destination)
            _update(callback, 80, "translation", "Excel 中没有需要翻译的中文自然语言文本，已原样交付")
            return 0

        _update(callback, 24, "translation", f"已建立稳定文本映射：{total} 条唯一文本；正在查询翻译记忆库")
        bilingual_target = getattr(client, "target_language_code", "") in {"zh-vi", "zh-en"}
        final_code = "vi" if getattr(client, "target_language_code", "") == "zh-vi" else "en"
        mapping: dict[str, str] = {}
        pending_sources: list[str] = []
        translation_inputs: list[str] = []

        if bilingual_target:
            for src in unique_texts:
                existing = _split_existing_bilingual_text(src)
                if existing:
                    mapping[src] = _normalize_bilingual_value(existing[0], existing[1])
                    continue
                unit = _extract_chinese_translation_unit(src)
                if not unit:
                    continue
                pending_sources.append(src)
                translation_inputs.append(unit)
            if translation_inputs:
                working_client = TranslationClient(
                    source_language=getattr(client, "source_language_code", "auto"),
                    target_language=final_code,
                )
                translated = working_client.translate_many(translation_inputs)
                for src, dst in zip(pending_sources, translated):
                    mapping[src] = _normalize_bilingual_value(src, str(dst or ""))
        else:
            translated = client.translate_many(unique_texts)
            for src, dst in zip(unique_texts, translated):
                value = str(dst or "").strip()
                if value:
                    mapping[src] = value
        changed = 0
        for _, _, nodes, original in text_nodes:
            value = mapping.get(original, original)
            if value == original:
                continue
            nodes[0].text = value
            nodes[0].set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
            for node in nodes[1:]:
                node.text = ""
            changed += 1

        _update(callback, 76, "translation", f"翻译完成：{changed} 个文本项；正在原位写回并校验包结构")
        with zipfile.ZipFile(destination, "w", compression=zipfile.ZIP_DEFLATED, compresslevel=6) as output:
            for info in archive.infolist():
                root = xml_payloads.get(info.filename + ":root")
                if isinstance(root, ET.Element):
                    data = ET.tostring(root, encoding="utf-8", xml_declaration=True)
                    output.writestr(info, data)
                else:
                    output.writestr(info, archive.read(info.filename))

    _apply_excel_multiline_layout(destination)
    with zipfile.ZipFile(source, "r") as before, zipfile.ZipFile(destination, "r") as after:
        missing = set(before.namelist()) - set(after.namelist())
        if missing:
            destination.unlink(missing_ok=True)
            raise RuntimeError(f"Excel 输出包缺少文件：{sorted(missing)[:5]}")
        if after.testzip() is not None:
            destination.unlink(missing_ok=True)
            raise RuntimeError("Excel 输出包完整性校验失败")
    _update(callback, 80, "translation", f"Excel 原位翻译完成；PLC 编号、中文和译文已分层排版；缓存命中 {getattr(client, 'persistent_cache_hits', 0)}；约剩余 0 秒")
    return changed


def _should_translate_excel_text(text: str) -> bool:
    value = str(text or "").strip()
    if not value or value.startswith("="):
        return False
    if value.startswith(("http://", "https://", "mailto:")):
        return False
    if value.upper() in {"#N/A", "#VALUE!", "#REF!", "#DIV/0!", "#NAME?", "#NUM!", "#NULL!"}:
        return False
    if re.fullmatch(r"[\d\s.,:;/%+\-_=()\[\]{}<>#@|\\]+", value):
        return False
    if re.fullmatch(r"(?:[A-Za-z]{1,8}\d+[A-Za-z0-9_.:/\-]*|[A-Za-z0-9]+(?:_[A-Za-z0-9]+)+|[A-Z]{2,}[A-Z0-9_.:/\-]*)", value):
        return False
    has_cjk = bool(_CJK_RE.search(value))
    has_latin_language = bool(re.search(r"[A-Za-zÀ-ỹ]{3,}", value))
    return has_cjk or has_latin_language


def _translate_pdf_to_docx(source: Path, destination: Path, client: TranslationClient, callback: ProgressCallback | None) -> int:
    reader = PdfReader(source)
    document = Document()
    total = max(1, len(reader.pages))
    translated = 0
    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            document.add_heading(f"Page {index}", level=2)
            for block in [part.strip() for part in text.split("\n\n") if part.strip()]:
                document.add_paragraph(client.translate(block))
                translated += 1
        _update(callback, 35 + int(index / total * 45), "translation", f"Translated PDF page {index}/{total}")
    document.save(destination)
    return translated




IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff", ".webp"}

def _translated_ocr_text(text: str, client: TranslationClient | None) -> str:
    if not text.strip() or client is None:
        return text
    # Translate line-by-line to retain invoice/list structure and reduce provider failures.
    output: list[str] = []
    for line in text.splitlines():
        if not line.strip():
            output.append("")
            continue
        translated = client.translate(line)
        output.append(str(translated or line).strip())
    return "\n".join(output).strip()

def _add_ocr_text(document: Document, text: str, target_code: str = "") -> int:
    count = 0
    for line in text.splitlines():
        if not line.strip():
            continue
        paragraph = document.add_paragraph()
        run = paragraph.add_run(line.strip())
        _apply_target_font(run, target_code)
        count += 1
    return count


def _detect_ocr_language(text: str) -> str:
    value = text or ""
    cjk = sum(1 for ch in value if "\u4e00" <= ch <= "\u9fff")
    vietnamese_marks = sum(1 for ch in value.lower() if ch in "ăâđêôơưáàảãạấầẩẫậắằẳẵặéèẻẽẹếềểễệíìỉĩịóòỏõọốồổỗộớờởỡợúùủũụứừửữựýỳỷỹỵ")
    latin = sum(1 for ch in value if ch.isascii() and ch.isalpha())
    if cjk and latin:
        return "中英混合"
    if cjk:
        return "中文"
    if vietnamese_marks:
        return "越南语"
    if latin:
        return "英文"
    return "未知"


def _looks_like_invoice(text: str) -> bool:
    upper = (text or "").upper()
    signals = ("INVOICE", "发票", "TOTAL", "总计", "UNIT PRICE", "单价", "AMOUNT", "金额")
    return sum(1 for token in signals if token in upper or token in text) >= 3


def _invoice_fields(lines: list[str]) -> tuple[list[str], list[list[str]], list[str]]:
    """Extract a practical invoice header, line-item table, and footer.

    This is intentionally conservative: it preserves every unparsed line in the
    header/footer rather than silently dropping OCR content.
    """
    clean = [" ".join(line.split()) for line in lines if line.strip()]
    header: list[str] = []
    rows: list[list[str]] = []
    footer: list[str] = []
    table_started = False
    for line in clean:
        normalized = line.replace("，", ",")
        if any(key in normalized.lower() for key in ("description", "unit price", "amount")) or all(key in normalized for key in ("数量", "金额")):
            table_started = True
            continue
        if any(key in normalized.lower() for key in ("total:", "payment:")) or normalized.startswith(("总计", "付款")):
            footer.append(normalized)
            table_started = False
            continue
        if table_started:
            # Typical invoice row: description + qty + unit price + amount.
            match = re.match(r"^(.*?)\s+(\d+)\s+([\d,.]+)\s+([\d,.]+(?:\s*[A-Za-z]{3})?)$", normalized)
            if match:
                rows.append([match.group(1).strip(), match.group(2), match.group(3), match.group(4)])
                continue
        if rows:
            footer.append(normalized)
        else:
            header.append(normalized)
    return header, rows, footer


def _set_cell_text(cell, value: str, target_code: str, bold: bool = False) -> None:
    cell.text = ""
    paragraph = cell.paragraphs[0]
    run = paragraph.add_run(value)
    run.bold = bold
    _apply_target_font(run, target_code)


def _build_editable_ocr_document(source: Path, translated: str, ocr_text: str, client: TranslationClient | None) -> tuple[Document, int, str]:
    """Build an editable, single-flow Word document from OCR results.

    V14 placed a generic title, the source image, and recognized text on three
    separate pages. V15 instead puts the editable result first and uses document-
    specific structures such as invoice tables. The original scan is appended as
    evidence without forcing an empty cover page.
    """
    document = Document()
    target_code = getattr(client, "target_language_code", "") if client else ""
    _configure_document_fonts(document, target_code)
    for section in document.sections:
        section.top_margin = Cm(1.7)
        section.bottom_margin = Cm(1.7)
        section.left_margin = Cm(1.8)
        section.right_margin = Cm(1.8)

    lines = [line for line in translated.splitlines() if line.strip()]
    blocks = 0
    template = "generic"
    if _looks_like_invoice(translated):
        template = "invoice"
        header, rows, footer = _invoice_fields(lines)
        title = next((line for line in header if "发票" in line or "INVOICE" in line.upper()), "发票识别结果")
        heading = document.add_heading(title, level=1)
        heading.alignment = 1
        blocks += 1
        for line in header:
            if line == title:
                continue
            p = document.add_paragraph()
            run = p.add_run(line)
            _apply_target_font(run, target_code)
            blocks += 1
        if rows:
            table = document.add_table(rows=1, cols=4)
            table.style = "Table Grid"
            headers = ["描述", "数量", "单价", "金额"] if target_code.startswith("zh") else ["Description", "Qty", "Unit Price", "Amount"]
            for idx, value in enumerate(headers):
                _set_cell_text(table.rows[0].cells[idx], value, target_code, bold=True)
            for row in rows:
                cells = table.add_row().cells
                for idx, value in enumerate(row):
                    _set_cell_text(cells[idx], value, target_code)
            blocks += len(rows) + 1
        for line in footer:
            p = document.add_paragraph()
            run = p.add_run(line)
            run.bold = line.startswith(("总计", "Total"))
            _apply_target_font(run, target_code)
            blocks += 1
    else:
        title = lines[0] if lines else ("OCR 识别结果" if target_code.startswith("zh") else "OCR Result")
        document.add_heading(title, level=1)
        blocks += 1
        for line in lines[1:]:
            p = document.add_paragraph()
            run = p.add_run(line)
            _apply_target_font(run, target_code)
            blocks += 1

    # Keep the scan as an appendix/evidence, not as the first user-visible page.
    if os.getenv("OCR_INCLUDE_SOURCE_IMAGE", "1").strip().lower() not in {"0", "false", "no"}:
        document.add_paragraph()
        label = "原始扫描件" if target_code.startswith("zh") else "Original scan"
        run = document.add_paragraph().add_run(label)
        run.bold = True
        _apply_target_font(run, target_code)
        try:
            document.add_picture(str(source), width=Inches(5.8))
        except Exception:
            pass

    document.core_properties.title = lines[0] if lines else source.stem
    document.core_properties.comments = "Editable OCR reconstruction by Document Automation AI V15"
    return document, blocks, template


def _image_to_docx(source: Path, destination: Path, client: TranslationClient | None, callback: ProgressCallback | None) -> dict[str, Any]:
    _update(callback, 38, "ocr", f"Reading image: {source.name}")
    preferred = getattr(client, "source_language_code", "auto") if client else "auto"
    ocr_text = extract_text_from_image(source, preferred)
    if not ocr_text.strip():
        raise RuntimeError(f"OCR produced no readable text for {source.name}.")
    _update(callback, 58, "ocr", f"OCR extracted {len(ocr_text)} characters")
    translated = _translated_ocr_text(ocr_text, client)
    document, blocks, template = _build_editable_ocr_document(source, translated, ocr_text, client)
    document.save(destination)
    validation = _validate_docx(destination)
    if validation["text_blocks"] == 0 or validation["characters"] == 0:
        raise RuntimeError("Generated OCR DOCX is empty; delivery was blocked.")
    return {
        "translated_items": blocks if client else 0,
        "ocr_characters": len(ocr_text),
        "ocr_text_blocks": len([x for x in ocr_text.splitlines() if x.strip()]),
        "recognized_language": _detect_ocr_language(ocr_text),
        "output_template": template,
        "editable_output": True,
        "validation": validation,
        "mode": "ocr_layout_reconstruction" if client else "ocr_layout",
    }

def _scanned_pdf_to_docx(source: Path, destination: Path, client: TranslationClient | None, callback: ProgressCallback | None) -> dict[str, Any]:
    if fitz is None:
        raise RuntimeError("PyMuPDF is required for scanned PDF OCR.")
    pdf = fitz.open(source)
    document = Document()
    target_code = getattr(client, "target_language_code", "") if client else ""
    _configure_document_fonts(document, target_code)
    total_chars = 0
    total_blocks = 0
    import tempfile
    with tempfile.TemporaryDirectory(prefix="docai_ocr_") as temp_dir:
        for index, page in enumerate(pdf, start=1):
            pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0), alpha=False)
            page_image = Path(temp_dir) / f"page_{index}.png"
            pix.save(str(page_image))
            preferred = getattr(client, "source_language_code", "auto") if client else "auto"
            ocr_text = extract_text_from_image(page_image, preferred)
            if not ocr_text.strip():
                continue
            translated = _translated_ocr_text(ocr_text, client)
            if index > 1:
                document.add_page_break()
            document.add_heading(f"Page {index}", level=1)
            total_blocks += _add_ocr_text(document, translated, target_code)
            if os.getenv("OCR_INCLUDE_SOURCE_IMAGE", "1").strip().lower() not in {"0", "false", "no"}:
                label = "原始扫描页" if target_code.startswith("zh") else "Original scanned page"
                run = document.add_paragraph().add_run(label)
                run.bold = True
                _apply_target_font(run, target_code)
                try:
                    document.add_picture(str(page_image), width=Inches(5.8))
                except Exception:
                    pass
            total_chars += len(ocr_text)
            _update(callback, 35 + int(index / max(1, len(pdf)) * 45), "ocr", f"OCR PDF page {index}/{len(pdf)}")
    pdf.close()
    if total_chars == 0:
        raise RuntimeError(f"OCR produced no readable text for scanned PDF {source.name}.")
    document.core_properties.comments = "Scanned PDF OCR and editable reconstruction by Document Automation AI V15"
    document.save(destination)
    validation = _validate_docx(destination)
    if validation["characters"] == 0:
        raise RuntimeError("Generated scanned-PDF DOCX is empty; delivery was blocked.")
    return {"translated_items": total_blocks if client else 0, "ocr_characters": total_chars, "ocr_text_blocks": total_blocks, "validation": validation, "mode": "scanned_pdf_ocr"}



def _normalize_text_value(value: str) -> str:
    value = value.replace("\u00a0", " ").replace("\u3000", " ")
    value = re.sub(r"[ \t]+", " ", value)
    return value.strip()


def _column_number(cell_ref: str) -> int:
    letters = re.match(r"[A-Za-z]+", cell_ref or "")
    if not letters:
        return 0
    value = 0
    for char in letters.group(0).upper():
        value = value * 26 + (ord(char) - 64)
    return value


def _cell_position(cell_ref: str) -> tuple[int, int]:
    match = re.fullmatch(r"([A-Za-z]+)(\d+)", cell_ref or "")
    if not match:
        return 0, 0
    return int(match.group(2)), _column_number(match.group(1))


def _compact_worksheet_xml(raw: bytes) -> tuple[bytes, dict[str, int]]:
    """Trim pathological trailing empty formatting cells without moving data.

    This implementation is deliberately row-oriented. A customer workbook may
    contain more than 18 million self-closing empty ``<c .../>`` elements; a
    cell-by-cell Python callback would be unacceptably slow. Meaningful cells are
    located from their closing tags, then each row is cut once at the first
    column beyond the protected data region.
    """
    ref_re = re.compile(rb'\br="([A-Za-z]+\d+)"')
    meaningful: list[tuple[int, int]] = []
    search_from = 0
    meaningful_cells = 0
    while True:
        close_at = raw.find(b"</c>", search_from)
        if close_at < 0:
            break
        open_at = raw.rfind(b"<c", 0, close_at)
        header_end = raw.find(b">", open_at, close_at) if open_at >= 0 else -1
        if open_at >= 0 and header_end >= 0:
            ref_match = ref_re.search(raw[open_at:header_end + 1])
            if ref_match:
                meaningful.append(_cell_position(ref_match.group(1).decode("ascii", "ignore")))
                meaningful_cells += 1
        search_from = close_at + 4

    for merge_ref in re.findall(rb'<mergeCell\s+ref="([A-Za-z]+\d+):([A-Za-z]+\d+)"', raw):
        for ref in merge_ref:
            meaningful.append(_cell_position(ref.decode("ascii", "ignore")))

    if not meaningful:
        return raw, {"cells_seen": 0, "phantom_cells_removed": 0, "rows_removed": 0}

    max_row = max(row for row, _ in meaningful)
    max_col = max(col for _, col in meaningful)
    protected_row = max_row + 10
    protected_col = max_col + 5
    cutoff_column = get_column_letter(protected_col + 1)

    output = bytearray()
    cursor = 0
    removed_cells = 0
    rows_removed = 0
    rows_seen = 0
    while True:
        row_start = raw.find(b"<row", cursor)
        if row_start < 0:
            output.extend(raw[cursor:])
            break
        row_end = raw.find(b"</row>", row_start)
        if row_end < 0:
            output.extend(raw[cursor:])
            break
        row_end += len(b"</row>")
        output.extend(raw[cursor:row_start])
        row_block = raw[row_start:row_end]
        header_end = row_block.find(b">")
        row_match = re.search(rb'\br="(\d+)"', row_block[:header_end + 1])
        row_number = int(row_match.group(1)) if row_match else 0
        rows_seen += 1

        # Rows after the real data region that contain no meaningful cell are
        # pure formatting noise and can be removed as a whole.
        if row_number > protected_row and b"</c>" not in row_block:
            removed_cells += row_block.count(b"<c ") + row_block.count(b"<c>")
            rows_removed += 1
        else:
            cutoff_ref = f'r="{cutoff_column}{row_number}"'.encode("ascii")
            cutoff_at = row_block.find(cutoff_ref)
            if cutoff_at < 0:
                # Sparse worksheets may jump directly from A to XFD. Find the
                # first referenced cell beyond the protected column.
                for ref_match in ref_re.finditer(row_block):
                    ref_text = ref_match.group(1).decode("ascii", "ignore")
                    ref_row, ref_col = _cell_position(ref_text)
                    if ref_row == row_number and ref_col > protected_col:
                        cutoff_at = ref_match.start()
                        break
            if cutoff_at >= 0:
                cell_start = row_block.rfind(b"<c", 0, cutoff_at)
                if cell_start >= 0:
                    removed_cells += row_block[cell_start:].count(b"<c ") + row_block[cell_start:].count(b"<c>")
                    row_block = row_block[:cell_start] + b"</row>"
            output.extend(row_block)
        cursor = row_end

    compacted = bytes(output)
    end_ref = f"{get_column_letter(max(1, max_col))}{max(1, max_row)}".encode("ascii")
    compacted = re.sub(rb'<dimension\s+ref="[^"]+"\s*/>', b'<dimension ref="A1:' + end_ref + b'"/>', compacted, count=1)
    return compacted, {
        "cells_seen": meaningful_cells,
        "phantom_cells_removed": removed_cells,
        "rows_removed": rows_removed,
        "rows_seen": rows_seen,
        "used_rows": max_row,
        "used_columns": max_col,
    }


def _detect_header_rows(ws, max_scan: int = 20) -> tuple[int, int]:
    """Return (title_row, header_row) using conservative visual/data heuristics."""
    nonempty_rows: list[tuple[int, int, int]] = []
    max_col = min(max(1, ws.max_column), 200)
    for r in range(1, min(ws.max_row, max_scan) + 1):
        values = [ws.cell(r, c).value for c in range(1, max_col + 1)]
        nonempty = sum(v not in (None, '') for v in values)
        text = sum(isinstance(v, str) and bool(v.strip()) for v in values)
        if nonempty:
            nonempty_rows.append((r, nonempty, text))
    if not nonempty_rows:
        return 0, 0
    title_row = nonempty_rows[0][0]
    # A title is usually sparse; a header is the first denser textual row.
    header_row = title_row
    for r, nonempty, text in nonempty_rows:
        if r > title_row and nonempty >= 2 and text >= max(1, nonempty // 2):
            header_row = r
            break
    return title_row, header_row




def _copy_cell_visual(source_cell, target_cell) -> None:
    target_cell.value = source_cell.value
    if source_cell.has_style:
        target_cell._style = copy.copy(source_cell._style)
    if source_cell.number_format:
        target_cell.number_format = source_cell.number_format
    target_cell.font = copy.copy(source_cell.font)
    target_cell.fill = copy.copy(source_cell.fill)
    target_cell.border = copy.copy(source_cell.border)
    target_cell.alignment = copy.copy(source_cell.alignment)
    target_cell.protection = copy.copy(source_cell.protection)


def _worksheet_content_bounds(ws) -> tuple[int, int]:
    """Return the last row/column containing real content.

    Worksheet.max_column can be inflated by formatting-only cells. Print layout
    must be based on actual values/formulas, merged ranges with content, and
    drawing anchors rather than the formatted tail.
    """
    max_row = 1
    max_col = 1
    for row in ws.iter_rows():
        for cell in row:
            if cell.value not in (None, ""):
                max_row = max(max_row, cell.row)
                max_col = max(max_col, cell.column)
    for merged in ws.merged_cells.ranges:
        anchor = ws.cell(merged.min_row, merged.min_col).value
        if anchor not in (None, ""):
            max_row = max(max_row, merged.max_row)
            max_col = max(max_col, merged.max_col)
    for image in getattr(ws, "_images", []):
        anchor = getattr(image, "anchor", None)
        marker = getattr(anchor, "_to", None) or getattr(anchor, "from_", None)
        if marker is not None:
            max_row = max(max_row, int(getattr(marker, "row", 0)) + 1)
            max_col = max(max_col, int(getattr(marker, "col", 0)) + 1)
    return max_row, max_col

def _is_simple_list_sheet(ws, max_row: int, max_col: int) -> bool:
    if max_row < 80 or max_col > 4 or ws.merged_cells.ranges or getattr(ws, "_images", []):
        return False
    nonempty_rows = 0
    formula_cells = 0
    dense_rows = 0
    for row in ws.iter_rows(min_row=1, max_row=max_row, max_col=max_col):
        values = [cell.value for cell in row]
        occupied = sum(value not in (None, "") for value in values)
        if occupied:
            nonempty_rows += 1
        if occupied > 4:
            dense_rows += 1
        formula_cells += sum(isinstance(value, str) and value.startswith("=") for value in values)
    return nonempty_rows >= 40 and dense_rows == 0 and formula_cells == 0

def _reflow_sparse_tail_to_right(ws, max_row: int | None = None, max_col: int | None = None) -> dict[str, int]:
    """Reflow only a genuinely simple list into two A4 columns.

    Complex PLC matrices are never moved. The previous implementation tried to
    move a sparse tail inside a wide matrix and could create unreadable sheets.
    """
    max_row, max_col = (max_row or _worksheet_content_bounds(ws)[0], max_col or _worksheet_content_bounds(ws)[1])
    if not _is_simple_list_sheet(ws, max_row, max_col):
        return {"moved_rows": 0, "start_row": 0, "mode": "preserve_complex_layout"}
    rows = []
    for r in range(1, max_row + 1):
        if any(ws.cell(r, c).value not in (None, "") for c in range(1, max_col + 1)):
            rows.append(r)
    # Preserve likely title/header rows and split only the data body.
    title_row, header_row = _detect_header_rows(ws)
    first_data = max(header_row + 1, 1)
    data_rows = [r for r in rows if r >= first_data]
    if len(data_rows) < 60:
        return {"moved_rows": 0, "start_row": first_data, "mode": "short_list"}
    split = (len(data_rows) + 1) // 2
    left_rows = data_rows[:split]
    right_rows = data_rows[split:]
    offset = max_col + 2
    for src_r, dst_r in zip(right_rows, left_rows):
        for c in range(1, max_col + 1):
            src = ws.cell(src_r, c)
            dst = ws.cell(dst_r, c + offset)
            _copy_cell_visual(src, dst)
            src.value = None
        if ws.row_dimensions[src_r].height and not ws.row_dimensions[dst_r].height:
            ws.row_dimensions[dst_r].height = ws.row_dimensions[src_r].height
    for c in range(1, max_col + 1):
        src_letter = get_column_letter(c)
        dst_letter = get_column_letter(c + offset)
        ws.column_dimensions[dst_letter].width = ws.column_dimensions[src_letter].width
    return {"moved_rows": len(right_rows), "start_row": first_data, "mode": "simple_two_column", "output_max_col": max_col * 2 + 2}

def _set_horizontal_page_breaks(ws, max_col: int, width_budget: float = 95.0) -> int:
    """Add A4-width column breaks without shrinking a wide matrix to one page."""
    ws.col_breaks = type(ws.col_breaks)()
    pages = 1
    used = 0.0
    for c in range(1, max_col + 1):
        letter = get_column_letter(c)
        width = float(ws.column_dimensions[letter].width or 8.43)
        if c > 1 and used + width > width_budget:
            ws.col_breaks.append(Break(id=c - 1))
            pages += 1
            used = 0.0
        used += width
    return pages

def _organize_xlsx_a4(source: Path, destination: Path, callback: ProgressCallback | None = None) -> dict[str, Any]:
    """Create a readable enterprise-organized workbook.

    The organizer no longer forces A4 fitting, moves rows into artificial
    columns, or shrinks a complex PLC matrix. It preserves the worksheet
    structure and improves only visible readability: real content bounds,
    practical widths, wrapped bilingual text, readable row heights, and safe
    print metadata. Empty sheets remain empty but are never created or cleared.
    """
    _update(callback, 12, 'cleanup', '正在分析工作表结构、有效区域、列宽和文字可读性')
    prepared, temporary = _prepare_xlsx_for_processing(source, callback, 'cleanup')
    try:
        wb = load_workbook(prepared, data_only=False)
    except (KeyError, ValueError, OSError):
        with zipfile.ZipFile(source, 'r') as src_zip, zipfile.ZipFile(destination, 'w', compression=zipfile.ZIP_DEFLATED) as dst_zip:
            totals = {'mode': 'safe_compaction_fallback', 'sheets': 0, 'structure_changed': False}
            for info in src_zip.infolist():
                payload = src_zip.read(info.filename)
                if info.filename.startswith('xl/worksheets/') and info.filename.endswith('.xml'):
                    payload, stats = _compact_worksheet_xml(payload)
                    totals['sheets'] += 1
                    totals['structure_changed'] = totals['structure_changed'] or bool(stats.get('phantom_cells_removed') or stats.get('rows_removed'))
                dst_zip.writestr(info, payload)
        return totals

    totals: dict[str, Any] = {
        'mode': 'readable_enterprise_layout',
        'sheets': 0,
        'content_bounds': {},
        'empty_sheets': [],
        'structure_changed': False,
    }
    try:
        for index, ws in enumerate(wb.worksheets, start=1):
            max_row, max_col = _worksheet_content_bounds(ws)
            has_content = any(
                cell.value not in (None, '')
                for row in ws.iter_rows(min_row=1, max_row=max_row, max_col=max_col)
                for cell in row
            )
            if not has_content:
                totals['empty_sheets'].append(ws.title)
                totals['content_bounds'][ws.title] = {'max_row': 0, 'max_col': 0}
                continue

            title_row, header_row = _detect_header_rows(ws)
            totals['content_bounds'][ws.title] = {'max_row': max_row, 'max_col': max_col}

            # Never move cells. Remove historical page-fit settings that made
            # hundreds of columns unreadably small.
            ws.sheet_properties.pageSetUpPr.fitToPage = False
            ws.page_setup.fitToWidth = 0
            ws.page_setup.fitToHeight = 0
            ws.page_setup.scale = 100
            ws.print_area = f'A1:{get_column_letter(max_col)}{max_row}'
            ws.row_breaks = type(ws.row_breaks)()
            ws.col_breaks = type(ws.col_breaks)()
            ws.sheet_view.showGridLines = True
            ws.sheet_view.zoomScale = 100

            # Keep existing orientation when present; otherwise use landscape
            # for matrix sheets and portrait for narrow lists.
            if not ws.page_setup.orientation:
                ws.page_setup.orientation = 'landscape' if max_col > 8 else 'portrait'
            ws.page_margins.left = 0.3
            ws.page_margins.right = 0.3
            ws.page_margins.top = 0.45
            ws.page_margins.bottom = 0.45

            # Determine widths from visible content without exploding columns.
            for c in range(1, max_col + 1):
                letter = get_column_letter(c)
                existing = ws.column_dimensions[letter].width
                samples: list[str] = []
                for r in range(1, min(max_row, 400) + 1):
                    value = ws.cell(r, c).value
                    if value not in (None, ''):
                        samples.append(str(value))
                if not samples:
                    continue
                max_line = max((max((len(x) for x in text.splitlines()), default=0) for text in samples), default=0)
                bilingual = any(('\n' in text or ' —— ' in text) for text in samples)
                desired = max(8.5, min(22.0 if bilingual else 18.0, max_line * 1.15 + 2))
                if existing is None or existing < 6 or existing > 35:
                    ws.column_dimensions[letter].width = desired
                else:
                    ws.column_dimensions[letter].width = min(24.0, max(8.0, float(existing)))

            # Normalize readability while preserving colors, borders and values.
            for r in range(1, max_row + 1):
                line_count = 1
                char_count = 0
                row_has_text = False
                for c in range(1, max_col + 1):
                    cell = ws.cell(r, c)
                    value = cell.value
                    if value in (None, ''):
                        continue
                    if isinstance(value, str):
                        row_has_text = True
                        line_count = max(line_count, value.count('\n') + 1)
                        char_count = max(char_count, max((len(x) for x in value.splitlines()), default=0))
                        cell.alignment = cell.alignment.copy(
                            wrap_text=True,
                            vertical='center',
                            horizontal=cell.alignment.horizontal or 'left',
                        )
                        current_size = cell.font.sz or 11
                        if current_size < 10.5:
                            cell.font = cell.font.copy(size=10.5)
                    elif cell.alignment.vertical is None:
                        cell.alignment = cell.alignment.copy(vertical='center')

                if row_has_text:
                    required = max(20.0, 17.0 * line_count)
                    if char_count > 28:
                        required += min(24.0, ((char_count - 1) // 28) * 6.0)
                    existing_h = ws.row_dimensions[r].height or 0
                    ws.row_dimensions[r].height = min(90.0, max(existing_h, required))

            if title_row:
                for cell in ws[title_row][:max_col]:
                    if cell.value not in (None, ''):
                        cell.font = cell.font.copy(bold=True, size=max(13, cell.font.sz or 11))
                        cell.alignment = cell.alignment.copy(wrap_text=True, vertical='center')
            if header_row:
                for cell in ws[header_row][:max_col]:
                    if cell.value not in (None, ''):
                        cell.font = cell.font.copy(bold=True, size=max(10.5, cell.font.sz or 10.5))
                        cell.alignment = cell.alignment.copy(wrap_text=True, vertical='center', horizontal='center')
                ws.freeze_panes = ws.freeze_panes or f'A{header_row + 1}'
                ws.print_title_rows = f'{header_row}:{header_row}'

            totals['sheets'] += 1
            _update(callback, 18 + int(index / max(1, len(wb.worksheets)) * 70), 'cleanup', f'整理工作表 {index}/{len(wb.worksheets)}：{ws.title}，有效区域 {max_row} 行 × {max_col} 列')

        destination.parent.mkdir(parents=True, exist_ok=True)
        wb.save(destination)
    finally:
        wb.close()
        if temporary:
            temporary.unlink(missing_ok=True)

    with zipfile.ZipFile(destination, 'r') as archive:
        bad = archive.testzip()
        if bad:
            destination.unlink(missing_ok=True)
            raise RuntimeError(f'智能整理输出校验失败：{bad}')
    totals['source_size_bytes'] = source.stat().st_size
    totals['output_size_bytes'] = destination.stat().st_size
    _update(callback, 92, 'cleanup', f"智能整理完成：{totals['sheets']} 个有效工作表；未移动任何单元格")
    return totals



def _classify_plc_function(name: str) -> str:
    value = str(name or '').strip()
    rules = [
        ('安全与急停', ('急停', '安全门', '安全继电器')),
        ('操作按钮', ('启动', '暂停', '停止', '复位', '旋钮', '按钮')),
        ('载具输送', ('载具', '运输板', '送料', '收料', '皮带', '小车')),
        ('检测与视觉', ('检测', '视觉', '相机', '扫码', 'AOI')),
        ('气动与真空', ('气缸', '夹爪', '真空', '顶升', '阻挡')),
        ('报警与状态', ('报警', '故障', '状态', '到位', '感应')),
    ]
    for category, words in rules:
        if any(word in value for word in words):
            return category
    if '备用' in value:
        return '备用信号'
    return '其他'


def _split_code_name(value: object) -> tuple[str, str]:
    text = str(value or '').strip()
    if not text or text.upper() in {'#N/A', '#VALUE!', '#REF!'}:
        return '', ''
    text = re.sub(r'\s+', ' ', text)
    match = re.match(r'^([A-Za-z]{1,8}\d+|NOT\d+|DOG\d+|POT\d+)\s*[|：:]?\s*(.*)$', text, re.I)
    if match:
        return match.group(1).upper(), match.group(2).strip()
    return '', text


def _clean_reconstructed_value(value: object) -> str:
    text = str(value or '').strip()
    if not text or text.upper() in {'#N/A', '#VALUE!', '#REF!', '#DIV/0!', '#NAME?'}:
        return ''
    return re.sub(r'\s+', ' ', text).strip()


def _engineering_system(name: str) -> str:
    value = str(name or '')
    rules = [
        ('安全系统', ('急停', '安全门', '安全继电器', '光栅')),
        ('视觉检测系统', ('AOI', '视觉', '相机', '检测', '扫码')),
        ('气动系统', ('气缸', '夹爪', '真空', '吸盘', '顶升', '阻挡')),
        ('输送物流系统', ('载具', '运输板', '皮带', '送料', '收料', '上料', '下料')),
        ('运动控制系统', ('轴', '伺服', '电机', '升降机')),
        ('操作与报警', ('启动', '停止', '暂停', '复位', '报警', '故障', '提示')),
    ]
    for system, words in rules:
        if any(word in value for word in words):
            return system
    return '其他系统'



_AUTOMATION_ZH_VI = {
    '操作条件不满足':'Điều kiện thao tác chưa được đáp ứng','已配置':'Đã cấu hình','未配置':'Chưa cấu hình',
    '工位':'Trạm làm việc','子模块':'Mô-đun con','所属系统':'Hệ thống','功能分类':'Phân loại chức năng',
    '输送物流系统':'Hệ thống vận chuyển và logistics','气动系统':'Hệ thống khí nén',
    '视觉检测系统':'Hệ thống kiểm tra thị giác','运动控制系统':'Hệ thống điều khiển chuyển động',
    '安全系统':'Hệ thống an toàn','操作与报警':'Vận hành và cảnh báo','通用设备系统':'Hệ thống thiết bị chung',
    '安全与急停':'An toàn và dừng khẩn cấp','操作按钮':'Nút thao tác','载具输送':'Vận chuyển đồ gá',
    '检测与视觉':'Kiểm tra và thị giác','气动与真空':'Khí nén và chân không','报警与状态':'Cảnh báo và trạng thái',
    '备用信号':'Tín hiệu dự phòng','其他':'Khác','运动轴':'Trục chuyển động','气缸':'Xi lanh','真空':'Chân không',
    '感应器':'Cảm biến','相机':'Camera','皮带线':'Băng tải','AOI/视觉':'AOI / Thị giác','安全门':'Cửa an toàn',
    '风扇':'Quạt','扫码器':'Máy quét mã','压力传感器':'Cảm biến áp suất','位移传感器':'Cảm biến dịch chuyển',
    '电批':'Tua vít điện','启动':'Khởi động','暂停':'Tạm dừng','停止':'Dừng','复位':'Đặt lại',
    '上升':'Nâng lên','下降':'Hạ xuống','伸出':'Đưa ra','缩回':'Thu về','松开':'Nhả ra','夹紧':'Kẹp chặt',
    '打开':'Mở','关闭':'Đóng','原位':'Vị trí gốc','动位':'Vị trí tác động','到位':'Đã đến vị trí',
    '上层线体':'Dây chuyền tầng trên','下层线体':'Dây chuyền tầng dưới','拆卸轴':'Trục tháo',
    '视觉交互':'Tương tác thị giác','检测轴':'Trục kiểm tra','载具输送':'Vận chuyển đồ gá',
    '上下料模组':'Mô-đun nạp/xả liệu','料盘搬运':'Vận chuyển khay vật liệu','上料':'Nạp liệu','下料':'Xả liệu',
    '升降机':'Thang nâng','轴':'Trục','检测模组':'Mô-đun kiểm tra','已删除':'Đã xóa',
    '中文':'Tiếng Trung','越南语':'Tiếng Việt','名称':'Tên','状态':'Trạng thái','类别':'Loại','编号':'Mã số',
    '地址':'Địa chỉ','序号':'STT','提示组':'Nhóm gợi ý','重复次数':'Số lần lặp','原位置':'Vị trí gốc',
}


def _glossary_vi(text: object) -> str:
    value=_clean_reconstructed_value(text)
    if not value:
        return ''
    if value in _AUTOMATION_ZH_VI:
        return _AUTOMATION_ZH_VI[value]
    # Translate common composite engineering labels deterministically.
    replacements=sorted(_AUTOMATION_ZH_VI.items(), key=lambda x: len(x[0]), reverse=True)
    result=value
    changed=False
    for zh,vi in replacements:
        if zh in result:
            result=result.replace(zh,vi); changed=True
    return result if changed and not _CJK_RE.search(result) else ''


def _bilingual_header(zh: str, vi: str) -> str:
    return f'{zh}\n{vi}'


def _style_reconstructed_sheet(ws, widths: list[float], table_color: str = '1F4E78') -> None:
    header_fill = PatternFill('solid', fgColor=table_color)
    alt_fill = PatternFill('solid', fgColor='F3F7FB')
    thin = Side(style='thin', color='D9E2F3')
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    for cell in ws[1]:
        if cell.value not in (None, ''):
            cell.font = Font(name='Microsoft YaHei', size=10.5, bold=True, color='FFFFFF')
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)
            cell.border = border
    ws.row_dimensions[1].height = 44
    for r, row in enumerate(ws.iter_rows(min_row=2), start=2):
        has_value = False
        max_lines = 1
        for cell in row:
            if cell.value not in (None, ''):
                has_value = True
                max_lines = max(max_lines, str(cell.value).count('\n') + 1)
            cell.font = Font(name='Microsoft YaHei', size=10.5)
            cell.alignment = Alignment(vertical='center', wrap_text=True)
            cell.border = border
            if r % 2 == 0:
                cell.fill = alt_fill
        if has_value:
            ws.row_dimensions[r].height = max(26, min(78, 19 * max_lines))
    for idx, width in enumerate(widths, start=1):
        ws.column_dimensions[get_column_letter(idx)].width = width
    ws.freeze_panes = 'A2'
    ws.auto_filter.ref = ws.dimensions
    ws.sheet_view.showGridLines = False
    ws.sheet_view.zoomScale = 100
    ws.sheet_properties.pageSetUpPr.fitToPage = False
    ws.page_setup.fitToWidth = 0
    ws.page_setup.fitToHeight = 0
    ws.page_setup.scale = 100
    ws.page_setup.orientation = 'landscape'
    ws.page_margins.left = 0.25
    ws.page_margins.right = 0.25


def _reconstruct_plc_configuration(source: Path, destination: Path, callback: ProgressCallback | None = None) -> dict[str, Any] | None:
    """Rebuild PLC/HMI matrices into concise bilingual engineering tables."""
    try:
        wb = load_workbook(source, data_only=False, read_only=False)
    except Exception:
        return None
    try:
        if not wb.worksheets:
            return None
        ws = wb.worksheets[0]
        probe = {str(ws.cell(r, 1).value or '').strip() for r in range(1, min(ws.max_row, 180) + 1)}
        required = {'输入名称', '轴名称', '气缸名称', '工位名称'}
        if len(required & probe) < 3:
            return None

        _update(callback, 14, 'cleanup', '识别为 PLC/HMI 配置表，正在提取、归类并建立企业双语术语')
        from openpyxl import Workbook
        out = Workbook()
        overview = out.active
        overview.title = '文档概览'
        overview.append([_bilingual_header('项目','Hạng mục'), _bilingual_header('内容','Nội dung')])

        plc_rows=[]
        for r in range(4, min(ws.max_row, 35) + 1):
            for c in range(2, min(ws.max_column, 65) + 1):
                code, name = _split_code_name(ws.cell(r, c).value)
                name=_clean_reconstructed_value(name)
                if code.startswith('X') and name and name not in {'备用','0'}:
                    category=_classify_plc_function(name); system=_engineering_system(name)
                    plc_rows.append([code,name,_glossary_vi(name),category,_glossary_vi(category),system,_glossary_vi(system),f'{get_column_letter(c)}{r}'])

        equipment_rows=[]
        equipment_map={'轴名称':'运动轴','气缸名称':'气缸','真空名称':'真空','感应器名称':'感应器','相机名称':'相机','皮带线名称':'皮带线','AOI名称':'AOI/视觉','安全门名称':'安全门','风扇名称':'风扇','扫码器名称':'扫码器','压力传感器名称':'压力传感器','位移传感器名称':'位移传感器','电批名称':'电批'}
        for r in range(1,min(ws.max_row,100)+1):
            category=equipment_map.get(str(ws.cell(r,1).value or '').strip())
            if not category: continue
            seen=set()
            for c in range(2,min(ws.max_column,65)+1):
                code,name=_split_code_name(ws.cell(r,c).value); name=_clean_reconstructed_value(name)
                if not (code or name) or name in {'0','备用'}: continue
                key=(code,name)
                if key in seen: continue
                seen.add(key)
                display=name or code; system=_engineering_system(category+display)
                equipment_rows.append([category,_glossary_vi(category),code,display,_glossary_vi(display),system,_glossary_vi(system),'已配置',_glossary_vi('已配置')])

        cylinder_rows=[]
        for c in range(2,min(ws.max_column,65)+1):
            code,name=_split_code_name(ws.cell(37,c).value); name=_clean_reconstructed_value(name)
            if not (code or name): continue
            vals=[_clean_reconstructed_value(ws.cell(rr,c).value) for rr in range(50,56)]
            if not any(vals) and not name: continue
            system=_engineering_system(name+'气缸')
            cylinder_rows.append([code,name,_glossary_vi(name),vals[0],vals[1],vals[2],_glossary_vi(vals[2]),vals[3],_glossary_vi(vals[3]),vals[4],vals[5],system,_glossary_vi(system)])

        station_rows=[]
        current_code=current_name=''
        for r in range(86,min(ws.max_row,260)+1):
            code=_clean_reconstructed_value(ws.cell(r,1).value)
            name=_clean_reconstructed_value(ws.cell(r,2).value)
            if not code or not name or code=='工位名称': continue
            if re.fullmatch(r'工位\d+',code):
                if re.fullmatch(r'\d+#',name): continue
                current_code,current_name=code,name
            elif re.fullmatch(r'工位\d+-\d+',code):
                if not current_code or re.fullmatch(r'\d+-\d+#',name): continue
                system=_engineering_system(current_name+name)
                station_rows.append([current_code,current_name,_glossary_vi(current_name),code,name,_glossary_vi(name),system,_glossary_vi(system)])
        with_children={r[0] for r in station_rows}
        for r in range(86,min(ws.max_row,260)+1):
            code=_clean_reconstructed_value(ws.cell(r,1).value); name=_clean_reconstructed_value(ws.cell(r,2).value)
            if re.fullmatch(r'工位\d+',code or '') and code not in with_children and name and not re.fullmatch(r'\d+#',name):
                system=_engineering_system(name)
                station_rows.append([code,name,_glossary_vi(name),'','','',system,_glossary_vi(system)])

        tip_counter={}
        for r in range(68,min(ws.max_row,180)+1):
            key=_clean_reconstructed_value(ws.cell(r,1).value)
            if not key.endswith('_Title'): continue
            group=key[:-6]
            for c in range(2,min(ws.max_column,65)+1):
                value=_clean_reconstructed_value(ws.cell(r,c).value)
                if not value or value=='0': continue
                tip_counter[(group,value)]=tip_counter.get((group,value),0)+1
        tip_rows=[]
        for (g,msg),count in tip_counter.items():
            system=_engineering_system(msg)
            tip_rows.append([g,msg,_glossary_vi(msg),count,system,_glossary_vi(system)])

        counts={'PLC输入信号':len(plc_rows),'设备清单':len(equipment_rows),'气缸IO配置':len(cylinder_rows),'工位结构':len(station_rows),'操作提示':len(tip_rows)}
        overview_rows=[
            ('文档类型','PLC/HMI 自动化设备工程文档'),('整理方式','按信号、设备、气缸、工位和提示信息分类汇总'),
            ('无效数据处理','已删除 #N/A、空白占位、备用项和无意义编号项'),('双语版式','中文与越南语严格使用相邻独立列显示'),
        ]
        for left,right in overview_rows: overview.append([f'{left}\n{_glossary_vi(left)}',right])
        overview.append([_bilingual_header('分类表','Bảng phân loại'),_bilingual_header('记录数量','Số lượng bản ghi')])
        for name,count in counts.items():
            overview.append([name,count]); overview.cell(overview.max_row,1).hyperlink=f"#'{name}'!A1"; overview.cell(overview.max_row,1).style='Hyperlink'

        plc=out.create_sheet('PLC输入信号'); plc.append([_bilingual_header('序号','STT'),_bilingual_header('地址','Địa chỉ'),_bilingual_header('中文功能','Chức năng tiếng Trung'),_bilingual_header('越南语功能','Chức năng tiếng Việt'),_bilingual_header('功能分类','Phân loại chức năng'),_bilingual_header('越南语分类','Phân loại tiếng Việt'),_bilingual_header('所属系统','Hệ thống'),_bilingual_header('越南语系统','Hệ thống tiếng Việt'),_bilingual_header('原位置','Vị trí gốc')])
        for i,row in enumerate(plc_rows,1): plc.append([i,*row])
        equipment=out.create_sheet('设备清单'); equipment.append([_bilingual_header('中文类别','Loại tiếng Trung'),_bilingual_header('越南语类别','Loại tiếng Việt'),_bilingual_header('编号','Mã số'),_bilingual_header('中文名称','Tên tiếng Trung'),_bilingual_header('越南语名称','Tên tiếng Việt'),_bilingual_header('所属系统','Hệ thống'),_bilingual_header('越南语系统','Hệ thống tiếng Việt'),_bilingual_header('中文状态','Trạng thái tiếng Trung'),_bilingual_header('越南语状态','Trạng thái tiếng Việt')])
        for row in equipment_rows: equipment.append(row)
        cylinder=out.create_sheet('气缸IO配置'); cylinder.append([_bilingual_header('气缸编号','Mã xi lanh'),_bilingual_header('中文名称','Tên tiếng Trung'),_bilingual_header('越南语名称','Tên tiếng Việt'),_bilingual_header('原位输入','Đầu vào vị trí gốc'),_bilingual_header('动位输入','Đầu vào vị trí tác động'),_bilingual_header('中文原位动作','Thao tác gốc tiếng Trung'),_bilingual_header('越南语原位动作','Thao tác gốc tiếng Việt'),_bilingual_header('中文动位动作','Thao tác động tiếng Trung'),_bilingual_header('越南语动位动作','Thao tác động tiếng Việt'),_bilingual_header('原位输出','Đầu ra vị trí gốc'),_bilingual_header('动位输出','Đầu ra vị trí tác động'),_bilingual_header('所属系统','Hệ thống'),_bilingual_header('越南语系统','Hệ thống tiếng Việt')])
        for row in cylinder_rows: cylinder.append(row)
        stations=out.create_sheet('工位结构'); stations.append([_bilingual_header('工位编号','Mã trạm'),_bilingual_header('中文工位名称','Tên trạm tiếng Trung'),_bilingual_header('越南语工位名称','Tên trạm tiếng Việt'),_bilingual_header('子模块编号','Mã mô-đun con'),_bilingual_header('中文子模块名称','Tên mô-đun tiếng Trung'),_bilingual_header('越南语子模块名称','Tên mô-đun tiếng Việt'),_bilingual_header('所属系统','Hệ thống'),_bilingual_header('越南语系统','Hệ thống tiếng Việt')])
        for row in station_rows: stations.append(row)
        tips=out.create_sheet('操作提示'); tips.append([_bilingual_header('提示组','Nhóm gợi ý'),_bilingual_header('中文提示','Gợi ý tiếng Trung'),_bilingual_header('越南语提示','Gợi ý tiếng Việt'),_bilingual_header('重复次数','Số lần lặp'),_bilingual_header('所属系统','Hệ thống'),_bilingual_header('越南语系统','Hệ thống tiếng Việt')])
        for row in tip_rows: tips.append(row)

        _style_reconstructed_sheet(overview,[26,76],'1F4E78')
        _style_reconstructed_sheet(plc,[8,14,30,34,20,24,20,28,14],'2F75B5')
        _style_reconstructed_sheet(equipment,[18,22,14,30,34,20,28,16,20],'548235')
        _style_reconstructed_sheet(cylinder,[14,26,30,15,15,18,22,18,22,15,15,20,28],'8064A2')
        _style_reconstructed_sheet(stations,[14,26,30,16,28,32,20,28],'C55A11')
        _style_reconstructed_sheet(tips,[16,38,42,12,20,28],'BF9000')

        # Drop empty rows/columns and validate meaningful output before saving.
        for sheet in out.worksheets:
            while sheet.max_row > 1 and all(sheet.cell(sheet.max_row,c).value in (None,'') for c in range(1,sheet.max_column+1)):
                sheet.delete_rows(sheet.max_row)
            while sheet.max_column > 1 and all(sheet.cell(r,sheet.max_column).value in (None,'') for r in range(1,sheet.max_row+1)):
                sheet.delete_cols(sheet.max_column)

        destination.parent.mkdir(parents=True, exist_ok=True)
        out.save(destination); out.close()
        check=load_workbook(destination,read_only=True,data_only=False)
        try:
            actual={sh.title:max(0,sh.max_row-1) for sh in check.worksheets}
            if not any(actual.get(name,0) for name in counts):
                raise RuntimeError('企业重构结果为空，已阻止交付')
        finally: check.close()
        _update(callback,92,'cleanup',f"智能重构完成：5 张分类表；PLC {counts['PLC输入信号']} 条；术语库已预翻译固定动作、状态和系统")
        return {'mode':'enterprise_reconstruction','document_type':'plc_hmi_configuration','sheets':len(actual),'records':actual,'structure_changed':True,'summary_counts':counts,'terminology_version':'21.3'}
    finally:
        wb.close()

def _clean_xlsx(source: Path, destination: Path, callback: ProgressCallback | None = None) -> dict[str, Any]:
    reconstructed = _reconstruct_plc_configuration(source, destination, callback)
    if reconstructed is not None:
        return reconstructed
    return _organize_xlsx_a4(source, destination, callback)


def _create_enterprise_analysis(source: Path, destination: Path, callback: ProgressCallback | None = None) -> dict[str, Any]:
    workbook = load_workbook(source, data_only=False, read_only=True)
    report = load_workbook(source, read_only=True, data_only=True)
    # Build a standalone analysis workbook.
    from openpyxl import Workbook
    out = Workbook()
    summary = out.active
    summary.title = "企业数据分析"
    headers=["工作表","有效行数","有效列数","非空单元格","公式数量","空值数量","重复行数"]
    summary.append(headers)
    for c in summary[1]:
        c.font=Font(bold=True)
        c.fill=PatternFill("solid", fgColor="DCE6F1")
        c.alignment=Alignment(horizontal="center")
    total_sheets=max(1,len(workbook.worksheets))
    totals={"rows":0,"cells":0,"formulas":0,"blanks":0,"duplicates":0}
    for idx,(ws,values_ws) in enumerate(zip(workbook.worksheets, report.worksheets), start=1):
        nonempty=0; formulas=0; blanks=0; seen=set(); duplicates=0; effective_rows=0; effective_cols=0
        for row in ws.iter_rows():
            vals=[]; row_nonempty=False
            for cell in row:
                val=cell.value; vals.append(val)
                if val not in (None, ""):
                    nonempty+=1; row_nonempty=True; effective_cols=max(effective_cols,cell.column)
                    if isinstance(val,str) and val.startswith('='): formulas+=1
                else: blanks+=1
            if row_nonempty:
                effective_rows=max(effective_rows,row[0].row if row else 0)
                key=tuple(vals[:min(len(vals),256)])
                if key in seen: duplicates+=1
                else: seen.add(key)
        summary.append([ws.title,effective_rows,effective_cols,nonempty,formulas,blanks,duplicates])
        totals["rows"]+=effective_rows; totals["cells"]+=nonempty; totals["formulas"]+=formulas; totals["blanks"]+=blanks; totals["duplicates"]+=duplicates
        _update(callback, 20 + int(idx / total_sheets * 70), "analysis_report", f"分析工作表 {idx}/{total_sheets}：{ws.title}")
    summary.append([])
    summary.append(["合计",totals["rows"],"",totals["cells"],totals["formulas"],totals["blanks"],totals["duplicates"]])
    summary.freeze_panes="A2"
    for col,width in enumerate([24,14,14,16,14,14,14], start=1): summary.column_dimensions[get_column_letter(col)].width=width
    out.save(destination)
    workbook.close(); report.close()
    return totals


def _process_file(original_name: str, stored_path: str, output_dir: Path, client: TranslationClient | None, callback: ProgressCallback | None, use_ocr: bool = False, use_cleanup: bool = False) -> dict[str, Any]:
    source = Path(stored_path)
    suffix = source.suffix.lower()
    safe_stem = Path(original_name).stem
    if use_cleanup and suffix == ".xlsx":
        cleaned = output_dir / f"{safe_stem}_智能整理.xlsx"
        cleanup_stats = _clean_xlsx(source, cleaned, callback)
        source = cleaned
        suffix = source.suffix.lower()
    else:
        cleanup_stats = None
    if client is None and not (use_ocr and (suffix in IMAGE_SUFFIXES or suffix == ".pdf")):
        destination = output_dir / Path(source.name).name
        if source.resolve() != destination.resolve():
            shutil.copy2(source, destination)
        return {"name": destination.name, "path": str(destination), "translated_items": 0, "mode": "cleaned" if cleanup_stats else "copied", "cleanup_stats": cleanup_stats}

    # Continue translation/conversion from the cleaned workbook when requested.
    original_source_name = original_name
    suffix_label = _target_suffix(client)
    details: dict[str, Any] = {}
    if suffix in IMAGE_SUFFIXES and use_ocr:
        destination = output_dir / f"{safe_stem}_{suffix_label}.docx"
        details = _image_to_docx(source, destination, client, callback)
        count = int(details.get("translated_items") or details.get("ocr_text_blocks") or 0)
    elif suffix == ".pptx":
        destination = output_dir / f"{safe_stem}_{suffix_label}.pptx"
        count = _translate_pptx(source, destination, client, callback)
    elif suffix == ".docx":
        destination = output_dir / f"{safe_stem}_{suffix_label}.docx"
        details = _translate_docx(source, destination, client, callback)
        count = int(details.get("translated_items") or 0)
    elif suffix == ".xlsx":
        destination = output_dir / f"{safe_stem}_{suffix_label}.xlsx"
        count = _translate_xlsx(source, destination, client, callback)
    elif suffix == ".pdf":
        destination = output_dir / f"{safe_stem}_{suffix_label}.docx"
        if use_ocr:
            # First try native PDF text. If it is effectively empty, render pages and OCR them.
            reader = PdfReader(source)
            native_text = "\n".join((page.extract_text() or "") for page in reader.pages).strip()
            if len(native_text) < 20:
                details = _scanned_pdf_to_docx(source, destination, client, callback)
                count = int(details.get("translated_items") or details.get("ocr_text_blocks") or 0)
            else:
                count = _translate_pdf_to_docx(source, destination, client, callback)
        else:
            count = _translate_pdf_to_docx(source, destination, client, callback)
    else:
        destination = output_dir / Path(original_name).name
        shutil.copy2(source, destination)
        count = 0
    if cleanup_stats and source != Path(stored_path) and destination != source:
        try:
            source.unlink(missing_ok=True)
        except OSError:
            LOGGER.warning("Unable to remove intermediate safe-cleanup workbook: %s", source)
    return {"name": destination.name, "path": str(destination), "translated_items": count, "mode": "translated" if count else "copied", "cleanup_stats": cleanup_stats, **details}


def run_local_job(order: dict[str, Any], source_paths: list[tuple[str, str]], output_dir: Path, progress_callback: ProgressCallback | None = None) -> dict[str, Any]:
    """Run an order with bounded file-level concurrency.

    V13 keeps one translation client per worker so provider usage counters and
    retries remain isolated.  Single-file orders keep the original sequential
    behavior, while batch orders process up to ``BATCH_MAX_WORKERS`` files in
    parallel (default 3) to reduce total waiting time without overloading the PC
    or translation provider.
    """
    output_dir.mkdir(parents=True, exist_ok=True)
    plan = build_plan(order)
    services = order.get("services") or []
    ocr = ocr_capability()
    translation = translation_capability()
    blockers: list[str] = []
    if "translation" in services and not translation.configured:
        blockers.append(translation.message)
    if "ocr" in services and not ocr.available:
        blockers.append(ocr.message)

    _update(progress_callback, 3, "validate", "正在检查源文件完整性、大小与格式")
    if not source_paths:
        raise RuntimeError("No source files were found for this order.")
    total_bytes = 0
    format_counts: dict[str, int] = {}
    for original_name, stored_path in source_paths:
        path = Path(stored_path)
        if not path.exists() or not path.is_file():
            raise RuntimeError(f"Source file is missing: {original_name}")
        total_bytes += path.stat().st_size
        suffix = path.suffix.lower().lstrip(".") or "unknown"
        format_counts[suffix] = format_counts.get(suffix, 0) + 1
    format_summary = "、".join(f"{key.upper()} {value}个" for key, value in sorted(format_counts.items()))
    _update(
        progress_callback, 10, "validate",
        f"已校验 {len(source_paths)} 个源文件，共 {total_bytes / 1024 / 1024:.2f} MB；{format_summary}",
    )

    _update(progress_callback, 12, "analyze", "正在读取已完成的文档结构分析结果")
    analysis_data = order.get("ai_analysis") or {}
    analysis_files = analysis_data.get("files") or []
    sheet_count = 0
    formula_count = 0
    merged_count = 0
    for item in analysis_files:
        details = item.get("details") or {}
        sheet_count += int(details.get("sheet_count") or 0)
        formula_count += int(details.get("formula_count_sample") or 0)
        merged_count += int(details.get("merged_range_count") or 0)
    detected_languages = "、".join(analysis_data.get("detected_languages") or []) or "自动识别"
    category = analysis_data.get("document_category") or "文档"
    detail_parts = [f"类别：{category}", f"语言：{detected_languages}"]
    if sheet_count:
        detail_parts.append(f"工作表：{sheet_count}个")
    if formula_count:
        detail_parts.append(f"公式样本：{formula_count}个")
    if merged_count:
        detail_parts.append(f"合并区域：{merged_count}个")
    _update(progress_callback, 20, "analyze", "文档结构分析完成；" + "；".join(detail_parts))

    manifest = {
        "order_number": order["order_number"],
        "created_at": now(),
        "plan": plan,
        "ocr": ocr.__dict__,
        "translation": translation.__dict__,
        "blockers": blockers,
        "source_files": [name for name, _ in source_paths],
        "batch": {"file_count": len(source_paths)},
    }
    manifest_path = output_dir / "processing_manifest.json"
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")

    if blockers:
        _update(progress_callback, 25, "configuration", blockers[0])
        return {"state": "waiting_configuration", "progress": 25, "current_step": "configuration", "plan": plan, "blockers": blockers, "manifest_path": str(manifest_path), "outputs": []}

    translation_data = order.get("translation") or {}

    shared_client: TranslationClient | None = None

    def make_client() -> TranslationClient | None:
        nonlocal shared_client
        if "translation" not in services:
            return None
        # Translation orders run sequentially by default. Reuse one client across
        # all files so repeated terms in st02-st09 are served from the same cache
        # instead of being sent to the AI provider again for every workbook.
        if shared_client is None:
            shared_client = TranslationClient(
                source_language=translation_data.get("source_language", "auto"),
                target_language=translation_data.get("target_language", "en"),
                custom_source=translation_data.get("custom_source_language", ""),
                custom_target=translation_data.get("custom_target_language", ""),
            )
        return shared_client

    if "translation" in services:
        _update(progress_callback, 30, "translation", "AI 翻译引擎连接成功")

    total = len(source_paths)
    # Translation is network-bound and providers commonly throttle concurrent
    # requests. Running several Excel files at once can leave every worker
    # waiting on a provider timeout. Default translation batches to one worker;
    # conversion-only orders may still use the normal bounded concurrency.
    worker_env = "BATCH_TRANSLATION_MAX_WORKERS" if "translation" in services else "BATCH_MAX_WORKERS"
    worker_default = "1" if "translation" in services else "3"
    max_workers = max(1, min(int(os.getenv(worker_env, worker_default)), total, 6))
    progress_lock = Lock()
    completed_count = 0
    file_fractions: dict[int, float] = {index: 0.0 for index in range(total)}
    last_reported_progress = 31
    outputs_by_index: dict[int, list[dict[str, Any]]] = {}
    usage_rows: list[dict[str, Any]] = []

    conversion_data = order.get("conversion") or {}
    requested_formats = conversion_data.get("formats") if "conversion" in services else ["original"]
    if not isinstance(requested_formats, list) or not requested_formats:
        requested_formats = ["original"]

    def report_file_progress(index: int, original_name: str, inner_progress: int, inner_step: str, message: str) -> None:
        nonlocal last_reported_progress
        fraction = max(0.0, min(1.0, float(inner_progress) / 100.0))
        with progress_lock:
            file_fractions[index] = max(file_fractions.get(index, 0.0), fraction)
            aggregate = sum(file_fractions.values()) / max(1, total)
            mapped = 32 + int(aggregate * 53)
            mapped = max(last_reported_progress, min(85, mapped))
            last_reported_progress = mapped
        _update(progress_callback, mapped, inner_step, f"{original_name}：{message}")

    def worker(index: int, original_name: str, stored_path: str):
        started_at = time.monotonic()
        client = make_client()
        file_failures: list[dict[str, Any]] = []
        LOGGER.info("Batch worker started: index=%s file=%s", index, original_name)
        try:
            def file_progress(inner_progress: int, inner_step: str, message: str) -> None:
                # Processing/translation occupies the first 68%% of each file's
                # aggregate share. This works for both single and batch orders.
                normalized = int(max(0, min(100, inner_progress)) * 0.68)
                report_file_progress(index, original_name, normalized, inner_step, message)

            primary_step = (
                "translation" if "translation" in services else
                "ocr" if "ocr" in services else
                "cleanup" if "data_cleanup" in services else
                "conversion" if "conversion" in services else
                "analyze"
            )
            report_file_progress(index, original_name, 1, primary_step, "开始处理文件")
            primary = _process_file(
                original_name,
                stored_path,
                output_dir,
                client,
                file_progress,
                use_ocr=("ocr" in services),
                use_cleanup=("data_cleanup" in services),
            )
            primary_path = Path(primary["path"])
            source_suffix = Path(original_name).suffix.lower().lstrip('.')
            effective_formats = list(requested_formats)
            if 'original' in effective_formats and source_suffix == 'pdf':
                effective_formats = [('pdf' if item == 'original' else item) for item in effective_formats]
            def conversion_progress(inner_progress: int, message: str) -> None:
                # Conversion occupies the remaining 32%% of this file's share.
                normalized = 68 + int(max(0, min(100, inner_progress)) * 0.32)
                report_file_progress(index, original_name, normalized, "conversion", message)

            if "conversion" in services:
                report_file_progress(index, original_name, 69, "conversion", "正在检查是否需要转换文件类型")
                converted_paths, conversion_records = convert_outputs(
                    primary_path,
                    effective_formats,
                    output_dir,
                    progress_callback=conversion_progress,
                )
            else:
                # Translation/OCR/organization-only orders keep the processed file
                # directly. Do not create a fake conversion stage or complete the
                # translation step early by emitting a conversion event.
                converted_paths = [primary_path]
                conversion_records = [{
                    "format": "original", "status": "completed", "path": str(primary_path),
                    "message": "保持原始文件类型并保留处理后的版式",
                }]
            file_failures.extend([
                {"source_name": original_name, **item}
                for item in conversion_records if item.get("status") == "failed"
            ])
            local_outputs = []
            for converted in converted_paths:
                local_outputs.append({
                    **primary,
                    "name": converted.name,
                    "path": str(converted),
                    "output_format": converted.suffix.lower().lstrip("."),
                    "source_name": original_name,
                    "conversion_records": conversion_records,
                })
            if "enterprise_analysis" in services and Path(stored_path).suffix.lower() == ".xlsx":
                analysis_path = output_dir / f"{Path(original_name).stem}_企业数据分析.xlsx"
                analysis_stats = _create_enterprise_analysis(Path(stored_path), analysis_path, file_progress)
                local_outputs.append({
                    "name": analysis_path.name, "path": str(analysis_path), "output_format": "xlsx",
                    "source_name": original_name, "mode": "enterprise_analysis", "translated_items": 0,
                    "analysis_stats": analysis_stats, "conversion_records": []
                })
            final_file_step = "conversion" if "conversion" in services else primary_step
            report_file_progress(index, original_name, 100, final_file_step, "文件处理完成")
            LOGGER.info(
                "Batch worker completed: index=%s file=%s elapsed=%.2fs outputs=%s failures=%s",
                index, original_name, time.monotonic() - started_at, len(local_outputs), len(file_failures),
            )
            return index, local_outputs, file_failures, (client.usage_summary() if client is not None else None)
        except Exception as exc:
            failure_step = "conversion" if "conversion" in services else ("translation" if "translation" in services else "analyze")
            report_file_progress(index, original_name, 100, failure_step, f"文件失败，已隔离：{exc}")
            LOGGER.exception(
                "Batch worker failed: index=%s file=%s elapsed=%.2fs",
                index, original_name, time.monotonic() - started_at,
            )
            file_failures.append({
                "source_name": original_name, "format": "processing", "status": "failed", "error": str(exc)
            })
            return index, [], file_failures, (client.usage_summary() if client is not None else None)

    failure_rows: list[dict[str, Any]] = []
    if total == 1 or max_workers == 1:
        for index, (original_name, stored_path) in enumerate(source_paths):
            current_stage = "translation" if "translation" in services else ("ocr" if "ocr" in services else ("cleanup" if "data_cleanup" in services else "conversion"))
            _update(progress_callback, max(32, last_reported_progress), current_stage, f"当前文件 {index + 1}/{total}：{original_name}；正在准备处理")
            idx, output, failures, usage = worker(index, original_name, stored_path)
            outputs_by_index[idx] = output
            failure_rows.extend(failures)
            if usage:
                usage_rows.append(usage)
            completed_count += 1
            report_file_progress(index, original_name, 100, "conversion" if "conversion" in services else "translation", f"已完成文件 {completed_count}/{total}")
    else:
        manifest["batch"].update({"mode": "parallel", "max_workers": max_workers})
        _update(progress_callback, 32, "conversion" if "conversion" in services else "translation", f"并行批处理已启动：{total} 个文件，{max_workers} 个工作线程")
        with ThreadPoolExecutor(max_workers=max_workers, thread_name_prefix="doc-ai") as pool:
            futures = {pool.submit(worker, i, name, path): (i, name) for i, (name, path) in enumerate(source_paths)}
            for future in as_completed(futures):
                index, original_name = futures[future]
                try:
                    idx, output, failures, usage = future.result()
                except Exception as exc:
                    idx, output, failures, usage = index, [], [{"source_name": original_name, "format": "processing", "status": "failed", "error": str(exc)}], None
                outputs_by_index[idx] = output
                failure_rows.extend(failures)
                if usage:
                    usage_rows.append(usage)
                with progress_lock:
                    completed_count += 1
                report_file_progress(index, original_name, 100, "conversion" if "conversion" in services else "translation", f"已完成文件 {completed_count}/{total}")

    outputs = [item for i in range(total) for item in outputs_by_index.get(i, [])]
    aggregate_usage = None
    if usage_rows:
        aggregate_usage = {
            "input_tokens": sum(int(r.get("input_tokens") or 0) for r in usage_rows),
            "output_tokens": sum(int(r.get("output_tokens") or 0) for r in usage_rows),
            "total_tokens": sum(int(r.get("total_tokens") or 0) for r in usage_rows),
            "estimated_cost_usd": round(sum(float(r.get("estimated_cost_usd") or 0) for r in usage_rows), 6),
            "files": len(usage_rows),
            "memory_cache_hits": sum(int(r.get("memory_cache_hits") or 0) for r in usage_rows),
            "session_cache_hits": sum(int(r.get("session_cache_hits") or 0) for r in usage_rows),
        }

    successful_output_count = len(outputs)
    failure_count = len(failure_rows)
    state, completion_message = _resolve_job_outcome(
        successful_output_count, failure_count, "manual_review" in services
    )

    # Failure details remain structured in the manifest/API. Never put an
    # error_report.txt into the customer's deliverables.
    quality_message = (
        f"质量检查通过：{successful_output_count} 个文件可交付"
        if failure_count == 0
        else f"质量检查完成：{successful_output_count} 个文件可交付，{failure_count} 项失败"
    )
    _update(progress_callback, 88, "quality", quality_message)
    if successful_output_count > 0:
        _update(progress_callback, 95, "export", f"已准备 {successful_output_count} 个交付文件")

    manifest["requested_output_formats"] = requested_formats
    manifest["outputs"] = outputs
    manifest["translation_usage"] = aggregate_usage
    manifest["failures"] = failure_rows
    manifest["partial_success"] = state == "partial_completed"
    manifest["successful_output_count"] = successful_output_count
    manifest["failure_count"] = failure_count
    manifest["terminal_state"] = state
    manifest["completed_at"] = now()
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    _update(progress_callback, 100, state, completion_message)
    return {
        "state": state, "progress": 100, "current_step": state, "plan": plan, "blockers": [],
        "manifest_path": str(manifest_path), "outputs": outputs, "translation_usage": aggregate_usage,
        "failures": failure_rows, "partial_success": state == "partial_completed",
        "successful_output_count": successful_output_count,
        "failure_count": failure_count,
        "completion_message": completion_message,
    }
