from __future__ import annotations

import csv
import html
import json
import os
import shutil
import subprocess
import tempfile
import threading
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Any, Callable

from docx import Document
from openpyxl import load_workbook
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

try:
    import fitz
except Exception:  # pragma: no cover
    fitz = None

_OFFICE_LOCK = threading.Lock()
_OFFICE_TIMEOUT = max(15, int(os.getenv("OFFICE_CONVERSION_TIMEOUT_SECONDS", "45")))
ConversionProgress = Callable[[int, str], None]


def _progress(callback: ConversionProgress | None, progress: int, message: str) -> None:
    if callback is not None:
        callback(max(0, min(100, progress)), message)


def _office_binary() -> str | None:
    return shutil.which("libreoffice") or shutil.which("soffice")


def _unique(path: Path) -> Path:
    if not path.exists():
        return path
    for index in range(2, 1000):
        candidate = path.with_name(f"{path.stem}_{index}{path.suffix}")
        if not candidate.exists():
            return candidate
    raise RuntimeError(f"Cannot create a unique output name for {path.name}")


def _libreoffice_convert(source: Path, target_ext: str, output_dir: Path) -> Path:
    binary = _office_binary()
    if not binary:
        raise RuntimeError("LibreOffice is required for this format conversion. Please install LibreOffice.")
    output_dir.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="daai_lo_") as profile:
        command = [
            binary, "--headless", f"-env:UserInstallation=file://{Path(profile).as_posix()}",
            "--convert-to", target_ext.lstrip("."), "--outdir", str(output_dir), str(source),
        ]
        with _OFFICE_LOCK:
            result = subprocess.run(command, capture_output=True, text=True, timeout=_OFFICE_TIMEOUT)
    expected = output_dir / f"{source.stem}.{target_ext.lstrip('.')}"
    if result.returncode != 0 or not expected.exists():
        detail = (result.stderr or result.stdout or "conversion produced no file").strip()
        raise RuntimeError(f"LibreOffice conversion failed for {source.name}: {detail}")
    return expected


def _powerpoint_convert(source: Path, destination: Path, format_code: int) -> Path:
    """Optional desktop Office automation; disabled by default to prevent UI popups."""
    if os.getenv("ALLOW_DESKTOP_OFFICE_AUTOMATION", "0").strip().lower() not in {"1", "true", "yes"}:
        raise RuntimeError("Desktop Office automation is disabled; using a headless converter instead.")
    if not shutil.which("powershell") and not shutil.which("pwsh"):
        raise RuntimeError("PowerShell is unavailable.")
    shell = shutil.which("powershell") or shutil.which("pwsh")
    src = str(source.resolve()).replace("'", "''")
    dst = str(destination.resolve()).replace("'", "''")
    script = f"""
$ErrorActionPreference = 'Stop'
$ppt = New-Object -ComObject PowerPoint.Application
try {{
  $ppt.Visible = 0
  $presentation = $ppt.Presentations.Open('{src}', $true, $false, $false)
  try {{ $presentation.SaveAs('{dst}', {format_code}) }} finally {{ $presentation.Close() }}
}} finally {{ $ppt.Quit() }}
"""
    with _OFFICE_LOCK:
        result = subprocess.run([shell, "-NoProfile", "-NonInteractive", "-Command", script], capture_output=True, text=True, timeout=_OFFICE_TIMEOUT)
    if result.returncode != 0 or not destination.exists():
        detail = (result.stderr or result.stdout or "PowerPoint conversion produced no file").strip()
        raise RuntimeError(detail)
    return destination




def _text_lines_to_pdf(lines: list[str], destination: Path, title: str = "") -> Path:
    """Create a dependable text PDF without launching Office applications."""
    if fitz is None:
        raise RuntimeError("PyMuPDF is required for portable PDF conversion.")
    pdf = fitz.open()
    page = pdf.new_page(width=595, height=842)
    y = 42
    all_lines = ([title] if title else []) + [str(x) for x in lines]
    for raw in all_lines:
        parts = [raw[i:i + 88] for i in range(0, max(1, len(raw)), 88)] or [""]
        for line in parts:
            if y > 800:
                page = pdf.new_page(width=595, height=842)
                y = 42
            safe = line.encode("latin-1", errors="replace").decode("latin-1")
            page.insert_text((42, y), safe, fontsize=10, fontname="helv")
            y += 15
        y += 3
    pdf.save(destination)
    pdf.close()
    return destination


def _docx_to_pdf_fallback(source: Path, destination: Path) -> Path:
    doc = Document(source)
    lines: list[str] = []
    lines.extend(p.text for p in doc.paragraphs if p.text.strip())
    for table in doc.tables:
        for row in table.rows:
            lines.append(" | ".join(cell.text for cell in row.cells))
    return _text_lines_to_pdf(lines, destination, source.stem)


def _xlsx_to_pdf_fallback(source: Path, destination: Path) -> Path:
    wb = load_workbook(source, data_only=False, read_only=True)
    lines: list[str] = []
    for ws in wb.worksheets:
        lines.append(f"[{ws.title}]")
        for row in ws.iter_rows(values_only=True):
            lines.append(" | ".join("" if value is None else str(value) for value in row))
    wb.close()
    return _text_lines_to_pdf(lines, destination, source.stem)


def _docx_to_xlsx(source: Path, destination: Path) -> Path:
    from openpyxl import Workbook
    doc = Document(source)
    wb = Workbook()
    ws = wb.active
    ws.title = "Document"
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            ws.append([paragraph.text])
    for table_index, table in enumerate(doc.tables, start=1):
        ws.append([])
        ws.append([f"Table {table_index}"])
        for row in table.rows:
            ws.append([cell.text for cell in row.cells])
    wb.save(destination)
    return destination


def _docx_to_pptx(source: Path, destination: Path) -> Path:
    doc = Document(source)
    prs = Presentation()
    chunks = [p.text for p in doc.paragraphs if p.text.strip()]
    if not chunks:
        chunks = [source.stem]
    for index in range(0, len(chunks), 8):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = source.stem if index == 0 else f"{source.stem} ({index // 8 + 1})"
        slide.placeholders[1].text = "\n".join(chunks[index:index + 8])
    prs.save(destination)
    return destination

def _pptx_to_pdf_fallback(source: Path, destination: Path) -> Path:
    """Portable readable PDF fallback when Office/LibreOffice is unavailable."""
    if fitz is None:
        raise RuntimeError("PyMuPDF is required for the PPT PDF fallback.")
    prs = Presentation(source)
    pdf = fitz.open()
    for slide_no, slide in enumerate(prs.slides, start=1):
        page = pdf.new_page(width=842, height=595)
        blocks = [f"Slide {slide_no}"]
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text.strip())
                if text:
                    blocks.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    blocks.append(" | ".join(cell.text for cell in row.cells))
        text = "\n\n".join(blocks)[:12000]
        page.insert_textbox(fitz.Rect(42, 35, 800, 560), text, fontsize=12, fontname="helv", lineheight=1.25)
    pdf.save(destination)
    pdf.close()
    return destination


def _pptx_to_docx(source: Path, destination: Path) -> Path:
    prs = Presentation(source)
    doc = Document()
    doc.add_heading(source.stem, 0)
    for slide_no, slide in enumerate(prs.slides, start=1):
        doc.add_heading(f"Slide {slide_no}", level=1)
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text.strip())
                if text:
                    doc.add_paragraph(text)
            if getattr(shape, "has_table", False):
                table = shape.table
                out = doc.add_table(rows=len(table.rows), cols=len(table.columns))
                out.style = "Table Grid"
                for r, row in enumerate(table.rows):
                    for c, cell in enumerate(row.cells):
                        out.cell(r, c).text = cell.text
        if slide_no < len(prs.slides):
            doc.add_page_break()
    doc.save(destination)
    return destination


def _pptx_to_xlsx(source: Path, destination: Path) -> Path:
    """Export slide text and tables to a structured workbook without Office/WPS."""
    from openpyxl import Workbook
    prs = Presentation(source)
    wb = Workbook()
    default = wb.active
    wb.remove(default)
    for slide_no, slide in enumerate(prs.slides, start=1):
        ws = wb.create_sheet(title=f"Slide {slide_no}"[:31])
        ws.append(["Type", "Content"])
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text.strip())
                if text:
                    ws.append(["Text", text])
            if getattr(shape, "has_table", False):
                ws.append(["Table", ""])
                for row in shape.table.rows:
                    ws.append([cell.text for cell in row.cells])
        if ws.max_row == 1:
            ws.append(["Slide", f"Slide {slide_no} contains no extractable text or table data."])
    if not prs.slides:
        ws = wb.create_sheet(title="Presentation")
        ws.append(["No slides found"])
    wb.save(destination)
    return destination


def _xlsx_to_csv(source: Path, destination: Path) -> Path:
    wb = load_workbook(source, data_only=False, read_only=True)
    ws = wb[wb.sheetnames[0]]
    with destination.open("w", newline="", encoding="utf-8-sig") as handle:
        writer = csv.writer(handle)
        for row in ws.iter_rows(values_only=True):
            writer.writerow(["" if value is None else value for value in row])
    wb.close()
    return destination


def _csv_to_xlsx(source: Path, destination: Path) -> Path:
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    raw = source.read_bytes()
    text = None
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "cp1252"):
        try:
            text = raw.decode(encoding)
            break
        except UnicodeDecodeError:
            continue
    if text is None:
        text = raw.decode("utf-8", errors="replace")
    for row in csv.reader(text.splitlines()):
        ws.append(row)
    wb.save(destination)
    return destination


def _read_text_file(source: Path) -> str:
    raw = source.read_bytes()
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "cp1252"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace")


def _extract_portable_content(source: Path) -> tuple[list[str], list[list[Any]]]:
    """Extract a conservative text/table representation for portable exports."""
    suffix = source.suffix.lower()
    blocks: list[str] = []
    rows: list[list[Any]] = []
    if suffix == ".docx":
        document = Document(source)
        blocks.extend(p.text for p in document.paragraphs if p.text.strip())
        for table in document.tables:
            for row in table.rows:
                values = [cell.text for cell in row.cells]
                rows.append(values)
                blocks.append(" | ".join(values))
    elif suffix == ".xlsx":
        workbook = load_workbook(source, data_only=False, read_only=True)
        for sheet in workbook.worksheets:
            blocks.append(f"[{sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                values = ["" if value is None else value for value in row]
                rows.append(values)
        workbook.close()
    elif suffix == ".pptx":
        presentation = Presentation(source)
        for slide_no, slide in enumerate(presentation.slides, start=1):
            blocks.append(f"Slide {slide_no}")
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    value = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text.strip())
                    if value:
                        blocks.append(value)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        values = [cell.text for cell in row.cells]
                        rows.append(values)
                        blocks.append(" | ".join(values))
    elif suffix == ".pdf" and fitz is not None:
        pdf = fitz.open(source)
        for index, page in enumerate(pdf, start=1):
            blocks.append(f"Page {index}")
            text = page.get_text("text").strip()
            if text:
                blocks.append(text)
        pdf.close()
    elif suffix == ".csv":
        text = _read_text_file(source)
        rows = [list(row) for row in csv.reader(text.splitlines())]
    elif suffix in {".txt", ".md", ".markdown", ".html", ".json", ".xml"}:
        blocks = [_read_text_file(source)]
    else:
        blocks = [source.name]
    return blocks, rows


def _portable_export(source: Path, destination: Path, fmt: str) -> Path:
    blocks, rows = _extract_portable_content(source)
    normalized = "md" if fmt in {"md", "markdown"} else fmt
    if normalized == "csv":
        with destination.open("w", newline="", encoding="utf-8-sig") as handle:
            writer = csv.writer(handle)
            for row in rows or [[block] for block in blocks]:
                writer.writerow(row)
    elif normalized == "json":
        destination.write_text(json.dumps({"source": source.name, "blocks": blocks, "rows": rows}, ensure_ascii=False, indent=2, default=str), encoding="utf-8")
    elif normalized == "xml":
        root = ET.Element("document", {"source": source.name})
        for index, block in enumerate(blocks, start=1):
            ET.SubElement(root, "block", {"index": str(index)}).text = str(block)
        for row_index, row in enumerate(rows, start=1):
            row_node = ET.SubElement(root, "row", {"index": str(row_index)})
            for column_index, value in enumerate(row, start=1):
                ET.SubElement(row_node, "cell", {"column": str(column_index)}).text = str(value)
        ET.ElementTree(root).write(destination, encoding="utf-8", xml_declaration=True)
    elif normalized == "html":
        parts = ["<!doctype html><html><head><meta charset=\"utf-8\"><title>", html.escape(source.stem), "</title></head><body>"]
        parts.extend(f"<p>{html.escape(str(block)).replace(chr(10), '<br>')}</p>" for block in blocks)
        if rows:
            parts.append("<table border=\"1\" cellspacing=\"0\" cellpadding=\"6\">")
            parts.extend("<tr>" + "".join(f"<td>{html.escape(str(value))}</td>" for value in row) + "</tr>" for row in rows)
            parts.append("</table>")
        parts.append("</body></html>")
        destination.write_text("".join(parts), encoding="utf-8")
    elif normalized == "md":
        text = [f"# {source.stem}", ""]
        text.extend(str(block) for block in blocks)
        if rows:
            text.extend(["", "## Structured data", ""])
            text.extend(" | ".join(str(value) for value in row) for row in rows)
        destination.write_text("\n\n".join(text), encoding="utf-8")
    else:
        destination.write_text("\n\n".join(str(block) for block in blocks) or "\n".join(",".join(map(str, row)) for row in rows), encoding="utf-8")
    return destination


def _portable_to_docx(source: Path, destination: Path) -> Path:
    blocks, rows = _extract_portable_content(source)
    document = Document()
    document.add_heading(source.stem, 0)
    for block in blocks:
        document.add_paragraph(str(block))
    if rows:
        width = max(1, max(len(row) for row in rows))
        table = document.add_table(rows=0, cols=width)
        table.style = "Table Grid"
        for values in rows:
            cells = table.add_row().cells
            for index, value in enumerate(values):
                cells[index].text = str(value)
    document.save(destination)
    return destination


def _portable_to_xlsx(source: Path, destination: Path) -> Path:
    from openpyxl import Workbook
    blocks, rows = _extract_portable_content(source)
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Content"
    if rows:
        for row in rows:
            sheet.append(list(row))
    else:
        for block in blocks:
            sheet.append([str(block)])
    workbook.save(destination)
    return destination


def _portable_to_pptx(source: Path, destination: Path) -> Path:
    blocks, rows = _extract_portable_content(source)
    content = [str(block) for block in blocks]
    content.extend(" | ".join(str(value) for value in row) for row in rows)
    presentation = Presentation()
    for index in range(0, max(1, len(content)), 8):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = source.stem if index == 0 else f"{source.stem} ({index // 8 + 1})"
        slide.placeholders[1].text = "\n".join(content[index:index + 8]) or source.name
    presentation.save(destination)
    return destination


def _image_to_pdf(source: Path, destination: Path) -> Path:
    with Image.open(source) as image:
        image.convert("RGB").save(destination, "PDF", resolution=150.0)
    return destination


def _image_to_pptx(source: Path, destination: Path) -> Path:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(source), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(destination)
    return destination


def _pdf_to_images(source: Path, output_dir: Path, base_name: str) -> list[Path]:
    if fitz is None:
        raise RuntimeError("PyMuPDF is required to export PDF pages as images.")
    folder = output_dir / f"{base_name}_images"
    folder.mkdir(parents=True, exist_ok=True)
    pdf = fitz.open(source)
    outputs: list[Path] = []
    for index, page in enumerate(pdf, start=1):
        target = folder / f"page_{index:03d}.png"
        page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False).save(str(target))
        outputs.append(target)
    pdf.close()
    return outputs


def convert_outputs(primary: Path, requested_formats: list[str], output_dir: Path, progress_callback: ConversionProgress | None = None) -> tuple[list[Path], list[dict[str, Any]]]:
    """Create every selected format from the processed primary file.

    `original` means keep the processed primary file. Unsupported combinations
    are reported in the manifest instead of silently pretending they succeeded.
    """
    requested = list(dict.fromkeys(requested_formats or ["original"]))
    _progress(progress_callback, 2, f"准备转换：{primary.name}")
    outputs: list[Path] = []
    records: list[dict[str, Any]] = []
    primary_ext = primary.suffix.lower().lstrip(".")

    if "original" in requested or primary_ext in requested:
        outputs.append(primary)
        records.append({"format": "original", "status": "completed", "path": str(primary)})

    conversion_targets = [item for item in requested if item not in {"original", primary_ext}]
    if not conversion_targets:
        _progress(progress_callback, 100, f"无需格式转换，直接交付 {primary.name}")
        return outputs, records
    target_total = len(conversion_targets)
    target_index = 0

    for fmt in requested:
        if fmt == "original" or fmt == primary_ext:
            continue
        try:
            target_index += 1
            _progress(progress_callback, 8 + int((target_index - 1) / target_total * 82), f"正在转换为 {fmt.upper()}：{primary.name}")
            target: Path | None = None
            if fmt == "pdf":
                if primary.suffix.lower() == ".pptx":
                    target = _unique(output_dir / f"{primary.stem}.pdf")
                    # Never launch WPS/PowerPoint in the normal server path. Prefer
                    # LibreOffice headless and use the portable renderer as fallback.
                    try:
                        converted = _libreoffice_convert(primary, "pdf", output_dir)
                        if converted != target:
                            shutil.move(str(converted), str(target))
                    except Exception:
                        _pptx_to_pdf_fallback(primary, target)
                elif primary.suffix.lower() == ".docx":
                    target = _unique(output_dir / f"{primary.stem}.pdf")
                    try:
                        converted = _libreoffice_convert(primary, "pdf", output_dir)
                        if converted != target:
                            shutil.move(str(converted), str(target))
                    except Exception:
                        _docx_to_pdf_fallback(primary, target)
                elif primary.suffix.lower() == ".xlsx":
                    target = _unique(output_dir / f"{primary.stem}.pdf")
                    try:
                        converted = _libreoffice_convert(primary, "pdf", output_dir)
                        if converted != target:
                            shutil.move(str(converted), str(target))
                    except Exception:
                        _xlsx_to_pdf_fallback(primary, target)
                elif primary.suffix.lower() in {".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff", ".webp"}:
                    target = _image_to_pdf(primary, _unique(output_dir / f"{primary.stem}.pdf"))
                else:
                    blocks, rows = _extract_portable_content(primary)
                    lines = blocks + [" | ".join(str(value) for value in row) for row in rows]
                    target = _text_lines_to_pdf(lines, _unique(output_dir / f"{primary.stem}.pdf"), primary.stem)
            elif fmt == "docx":
                if primary.suffix.lower() == ".pptx":
                    target = _pptx_to_docx(primary, _unique(output_dir / f"{primary.stem}.docx"))
                elif primary.suffix.lower() in {".odt", ".rtf", ".html", ".txt"}:
                    try:
                        target = _libreoffice_convert(primary, "docx", output_dir)
                    except Exception:
                        target = _portable_to_docx(primary, _unique(output_dir / f"{primary.stem}.docx"))
                else:
                    target = _portable_to_docx(primary, _unique(output_dir / f"{primary.stem}.docx"))
            elif fmt == "xlsx":
                if primary.suffix.lower() == ".csv":
                    target = _csv_to_xlsx(primary, _unique(output_dir / f"{primary.stem}.xlsx"))
                elif primary.suffix.lower() == ".docx":
                    target = _docx_to_xlsx(primary, _unique(output_dir / f"{primary.stem}.xlsx"))
                elif primary.suffix.lower() == ".pptx":
                    target = _pptx_to_xlsx(primary, _unique(output_dir / f"{primary.stem}.xlsx"))
                else:
                    target = _portable_to_xlsx(primary, _unique(output_dir / f"{primary.stem}.xlsx"))
            elif fmt in {"txt", "md", "markdown", "html", "json", "xml", "csv"}:
                extension = "md" if fmt in {"md", "markdown"} else fmt
                target = _portable_export(primary, _unique(output_dir / f"{primary.stem}.{extension}"), fmt)
            elif fmt == "pptx":
                if primary.suffix.lower() in {".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff", ".webp"}:
                    target = _image_to_pptx(primary, _unique(output_dir / f"{primary.stem}.pptx"))
                elif primary.suffix.lower() == ".docx":
                    target = _docx_to_pptx(primary, _unique(output_dir / f"{primary.stem}.pptx"))
                else:
                    target = _portable_to_pptx(primary, _unique(output_dir / f"{primary.stem}.pptx"))
            elif fmt == "images":
                if primary.suffix.lower() == ".pdf":
                    image_paths = _pdf_to_images(primary, output_dir, primary.stem)
                    outputs.extend(image_paths)
                    records.append({"format": fmt, "status": "completed", "paths": [str(x) for x in image_paths]})
                    continue
                if primary.suffix.lower() in {".docx", ".xlsx", ".pptx"}:
                    pdf = _libreoffice_convert(primary, "pdf", output_dir)
                    image_paths = _pdf_to_images(pdf, output_dir, primary.stem)
                    outputs.extend(image_paths)
                    records.append({"format": fmt, "status": "completed", "paths": [str(x) for x in image_paths]})
                    continue
            if target is None or not target.exists():
                raise RuntimeError(f"Conversion {primary.suffix or 'file'} → {fmt} is not supported for this file.")
            outputs.append(target)
            records.append({"format": fmt, "status": "completed", "path": str(target)})
            _progress(progress_callback, 8 + int(target_index / target_total * 82), f"{fmt.upper()} 转换完成：{target.name}")
        except Exception as exc:
            records.append({"format": fmt, "status": "failed", "error": str(exc)})
            _progress(progress_callback, 8 + int(target_index / target_total * 82), f"{fmt.upper()} 转换失败，已隔离：{exc}")

    # Avoid duplicates produced by LibreOffice or overlapping selections.
    unique_outputs: list[Path] = []
    seen: set[str] = set()
    for item in outputs:
        key = str(item.resolve())
        if key not in seen and item.exists():
            seen.add(key)
            unique_outputs.append(item)
    _progress(progress_callback, 100, f"格式转换阶段完成：{len(unique_outputs)} 个输出文件")
    return unique_outputs, records
