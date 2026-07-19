from pathlib import Path
from datetime import datetime, timedelta
import csv, io, json, random, shutil, zipfile
from PIL import Image, ImageDraw, ImageFont, ImageFilter
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.section import WD_ORIENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from reportlab.lib.pagesizes import A4, landscape
from reportlab.pdfgen import canvas
from artifact_tool import Workbook, SpreadsheetFile

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / 'Enterprise_Test_Suite'
if OUT.exists(): shutil.rmtree(OUT)
for d in ['01_Word','02_Excel','03_PPT','04_PDF','05_ScanPDF','06_Image','07_CSV','08_ZIP','09_Error_Cases','10_Regression_Manifests']:
    (OUT/d).mkdir(parents=True, exist_ok=True)
random.seed(17)

def add_docx(name, kind, pages=14):
    doc=Document(); sec=doc.sections[0]; sec.top_margin=Inches(.55); sec.bottom_margin=Inches(.55)
    title=doc.add_heading(kind,0); title.alignment=WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph('Document Automation AI - Enterprise Acceptance Fixture | 中英越混合测试')
    for i in range(1,pages+1):
        doc.add_heading(f'{i}. Section / 章节 / Mục {i}', level=1)
        p=doc.add_paragraph(); p.add_run('Requirement: ').bold=True
        p.add_run('The automation cell shall maintain traceability, safety interlocks, and bilingual records. ' * 3)
        doc.add_paragraph(f'中文：设备编号 EQ-{i:03d}，验收标准包含精度、节拍、安全联锁与追溯。')
        doc.add_paragraph(f'Tiếng Việt: Thiết bị EQ-{i:03d}, tiêu chuẩn nghiệm thu gồm độ chính xác, chu kỳ và an toàn.')
        table=doc.add_table(rows=6, cols=5); table.style='Table Grid'
        headers=['Item','Specification','Measured','Status','Remark']
        for c,h in enumerate(headers): table.cell(0,c).text=h
        for r in range(1,6):
            vals=[f'{i}.{r}',f'Tolerance ±{r/10:.1f} mm',f'{10+i+r/10:.2f}','PASS' if r<5 else 'REVIEW','含特殊字符 Ω / µm / ≤']
            for c,v in enumerate(vals): table.cell(r,c).text=str(v)
        if i%3==0:
            doc.add_paragraph('NOTE: This page intentionally contains dense tables, multilingual text, symbols, and repeated headings.')
        doc.add_page_break()
    doc.sections[0].header.paragraphs[0].text=f'{kind} | CONFIDENTIAL | Revision V17'
    doc.sections[0].footer.paragraphs[0].text='Controlled enterprise acceptance sample'
    path=OUT/'01_Word'/name; doc.save(path)

for n,k,p in [
('01_Long_Equipment_Purchase_Contract.docx','Long Equipment Purchase Contract',22),
('02_FAT_SAT_Combined_Report.docx','FAT / SAT Combined Acceptance Report',18),
('03_Multilingual_Operations_Manual.docx','Multilingual Operations Manual',20),
('04_PLC_Software_Design_Specification.docx','PLC Software Design Specification',16),
('05_Robot_Cell_Safety_Risk_Assessment.docx','Robot Cell Safety Risk Assessment',15),
('06_Preventive_Maintenance_Handbook.docx','Preventive Maintenance Handbook',18),
('07_Supplier_Quality_Audit_Report.docx','Supplier Quality Audit Report',14),
('08_Change_Control_and_Deviation_Log.docx','Change Control and Deviation Log',14)]: add_docx(n,k,p)

def xlsx_file(name, rows, sheets):
    wb=Workbook.create()
    for sidx in range(sheets):
        sh=wb.worksheets.add(f'Sheet_{sidx+1}')
        sh.get_range('A1:J1').values=[['ID','Part No.','Description','Qty','Unit Price','Subtotal','Status','Owner','Due Date','Variance']]
        data=[]
        start=datetime(2026,1,1)
        for i in range(1,rows+1):
            data.append([i,f'PN-{sidx+1:02d}-{i:06d}',f'Automation component {i} / 零件 {i}',(i%17)+1,round(2.5+(i%83)*1.37,2),None,['Open','Approved','Hold'][i%3],['Lan','Quality','Engineering'][i%3],start+timedelta(days=i%365),None])
        sh.get_range(f'A2:J{rows+1}').values=data
        sh.get_range('F2').formulas=[['=D2*E2']]; sh.get_range(f'F2:F{rows+1}').fill_down()
        sh.get_range('J2').formulas=[['=F2-ROUND(F2*0.97,2)']]; sh.get_range(f'J2:J{rows+1}').fill_down()
        sh.freeze_panes.freeze_rows(1)
        sh.get_range('A1:J1').format={'fill':'#1F4E78','font':{'bold':True,'color':'#FFFFFF'},'wrap_text':True}
        sh.get_range(f'D2:F{rows+1}').format.number_format='#,##0.00'
        sh.get_range(f'I2:I{rows+1}').format.number_format='yyyy-mm-dd'
        sh.get_range(f'G2:G{rows+1}').data_validation={'rule':{'type':'list','values':['Open','Approved','Hold']}}
        sh.get_range(f'F2:F{rows+1}').conditional_formats.add_data_bar({'color':'#4472C4','gradient':True})
        sh.get_range('A:J').format.column_width=16; sh.get_range('C:C').format.column_width=30
        sh.tables.add(f'A1:J{rows+1}',True,f'Table{sidx+1}')
    SpreadsheetFile.export_xlsx(wb).save(str(OUT/'02_Excel'/name))

# Excel fixtures are generated separately with artifact_tool in build_v17_excel.py

def ppt(name, slides):
    prs=Presentation(); prs.slide_width=PInches(13.333); prs.slide_height=PInches(7.5)
    for i in range(slides):
        sl=prs.slides.add_slide(prs.slide_layouts[6])
        tb=sl.shapes.add_textbox(PInches(.6),PInches(.35),PInches(12),PInches(.7)); tf=tb.text_frame; tf.text=f'Enterprise Automation Review - Slide {i+1}'
        tf.paragraphs[0].font.size=PPt(26); tf.paragraphs[0].font.bold=True
        body=sl.shapes.add_textbox(PInches(.7),PInches(1.2),PInches(5.7),PInches(5.5)); bt=body.text_frame
        for j in range(8):
            p=bt.paragraphs[0] if j==0 else bt.add_paragraph(); p.text=f'{j+1}. KPI {i+1}-{j+1}: OEE {70+(i+j)%29}% | 中文说明 | Mô tả tiếng Việt'; p.font.size=PPt(15)
        chart=sl.shapes.add_chart(51,PInches(6.7),PInches(1.25),PInches(5.8),PInches(3.2),None) if False else None
        tbl=sl.shapes.add_table(6,4,PInches(6.7),PInches(4.7),PInches(5.8),PInches(2)).table
        for r in range(6):
            for c in range(4): tbl.cell(r,c).text=['Metric','Plan','Actual','Status'][c] if r==0 else f'{(i+1)*(r+c+1)}'
    prs.save(OUT/'03_PPT'/name)
for n,s in [('01_Automation_Solution_Proposal_40_Slides.pptx',40),('02_FAT_Management_Review_30_Slides.pptx',30),('03_Operator_Training_Deck_35_Slides.pptx',35),('04_Maintenance_Strategy_25_Slides.pptx',25)]: ppt(n,s)

def make_pdf(name,pages,land=False,two_col=False):
    ps=landscape(A4) if land else A4; c=canvas.Canvas(str(OUT/'04_PDF'/name),pagesize=ps); w,h=ps
    for p in range(1,pages+1):
        c.setFont('Helvetica-Bold',16); c.drawString(40,h-45,f'Enterprise Technical Report - Page {p}/{pages}')
        c.setFont('Helvetica',8)
        cols=[40,w/2+10] if two_col else [40]
        for x in cols:
            y=h-75
            for i in range(32):
                c.drawString(x,y,f'{p}.{i+1:02d} Requirement / acceptance result / tolerance +/-0.{i%9} mm / PASS')
                y-=14
        c.rect(40,55,w-80,90); c.drawString(50,125,'Complex table region: Item | Planned | Actual | Status | Remark')
        for i in range(5): c.line(40,55+i*18,w-40,55+i*18)
        c.showPage()
    c.save()
for n,p,l,t in [('01_Long_Technical_Contract_80_Pages.pdf',80,False,False),('02_Two_Column_Compliance_Report.pdf',45,False,True),('03_Landscape_Engineering_Drawings.pdf',30,True,False),('04_Mixed_Table_Text_Report.pdf',55,False,False),('05_Multilingual_Quality_Report.pdf',40,False,True)]: make_pdf(n,p,l,t)

def scan_fixture(base, idx, rotate=0, blur=0, shadow=False):
    im=Image.new('RGB',(1654,2339),'white'); d=ImageDraw.Draw(im)
    d.rectangle((90,80,1560,2250),outline='black',width=4)
    d.text((130,120),f'SCANNED ENTERPRISE DOCUMENT #{idx}',fill='black')
    y=210
    for r in range(45):
        d.text((130,y),f'Line {r+1:02d}: Equipment EQ-{idx:03d}; value={100+r/10:.1f}; status=PASS; OCR symbols 0/O 1/I 5/S',fill='black'); y+=42
    if shadow:
        sh=Image.new('RGBA',im.size,(0,0,0,0)); sd=ImageDraw.Draw(sh); sd.rectangle((0,0,420,2339),fill=(0,0,0,80)); im=Image.alpha_composite(im.convert('RGBA'),sh).convert('RGB')
    if blur: im=im.filter(ImageFilter.GaussianBlur(blur))
    if rotate: im=im.rotate(rotate,expand=True,fillcolor='white')
    im.save(base)
for i,(rot,blur,shadow) in enumerate([(2,0,False),(-4,1,False),(0,2,True),(7,1,True),(0,0,False),(3,2,False)],1):
    scan_fixture(OUT/'06_Image'/f'{i:02d}_OCR_Scan_Rotation_{rot}_Blur_{blur}.png',i,rot,blur,shadow)
# multipage scan PDFs
imgs=list((OUT/'06_Image').glob('*.png'))
for i in range(3):
    pdf_path=OUT/'05_ScanPDF'/f'{i+1:02d}_MultiPage_Scanned_Document.pdf'
    c=canvas.Canvas(str(pdf_path), pagesize=A4)
    w,h=A4
    for img_path in imgs[i:i+4]:
        c.drawImage(str(img_path),0,0,width=w,height=h,preserveAspectRatio=True,anchor='c')
        c.showPage()
    c.save()

for i in range(4):
    path=OUT/'07_CSV'/f'{i+1:02d}_Enterprise_Data_10000_Rows.csv'
    with path.open('w',newline='',encoding='utf-8-sig') as f:
        w=csv.writer(f); w.writerow(['ID','Equipment','Timestamp','Value','Unit','Alarm','Chinese','Vietnamese'])
        for r in range(10000): w.writerow([r+1,f'EQ-{r%250:03d}',f'2026-07-{1+r%28:02d} {r%24:02d}:{r%60:02d}',round(20+(r%100)/7,3),'mm','ALARM' if r%997==0 else '',f'设备{r%250}',f'Thiết bị {r%250}'])

# ZIP fixtures
for i in range(3):
    zp=OUT/'08_ZIP'/f'{i+1:02d}_Nested_Enterprise_Package.zip'
    with zipfile.ZipFile(zp,'w',zipfile.ZIP_DEFLATED) as z:
        z.writestr('中文目录/README.txt','Nested ZIP fixture for enterprise regression')
        z.writestr('level1/level2/level3/data.csv','a,b,c\n1,2,3\n')
        z.write(OUT/'01_Word'/'02_FAT_SAT_Combined_Report.docx','documents/FAT_Report.docx')

# Errors that do not break Windows extraction
(OUT/'09_Error_Cases'/'01_empty.pdf').write_bytes(b'')
(OUT/'09_Error_Cases'/'02_fake_pdf.pdf').write_text('This is not a PDF',encoding='utf-8')
(OUT/'09_Error_Cases'/'03_corrupt_docx.docx').write_bytes(b'PK\x03\x04BROKEN')
(OUT/'09_Error_Cases'/'04_corrupt_xlsx.xlsx').write_bytes(b'PK\x03\x04BROKEN')
(OUT/'09_Error_Cases'/'05_no_extension').write_text('no extension test',encoding='utf-8')
(OUT/'09_Error_Cases'/'06_special_chars_中文_#_&_test.txt').write_text('special filename',encoding='utf-8')

# Runtime long path generator avoids extraction failure
(OUT/'09_Error_Cases'/'generate_runtime_long_path.py').write_text("""from pathlib import Path\nroot=Path('runtime_long_path')\nname='segment_'+'x'*45\np=root\nfor _ in range(6): p=p/name\np.mkdir(parents=True,exist_ok=True)\n(p/'long_path_test.pdf').write_bytes(b'%PDF-1.4\\n% runtime long path fixture')\nprint(p)\n""",encoding='utf-8')

files=[p for p in OUT.rglob('*') if p.is_file()]
manifest=[]
for p in files:
    manifest.append({'path':str(p.relative_to(OUT)).replace('\\','/'),'size_bytes':p.stat().st_size,'category':p.parts[-2]})
(OUT/'10_Regression_Manifests'/'suite_manifest.json').write_text(json.dumps(manifest,ensure_ascii=False,indent=2),encoding='utf-8')
(OUT/'10_Regression_Manifests'/'49_file_smoke_test.txt').write_text('\n'.join(x['path'] for x in manifest[:49]),encoding='utf-8')
(OUT/'10_Regression_Manifests'/'100_file_regression.txt').write_text('\n'.join((x['path'] for x in (manifest*3)[:100])),encoding='utf-8')
(OUT/'10_Regression_Manifests'/'300_file_pressure_plan.txt').write_text('Use the generator script to duplicate files at runtime; do not store 300 copies in source control.\n',encoding='utf-8')
readme=f'''# V17 Enterprise Test Suite\n\nGenerated: {datetime.now().isoformat(timespec='seconds')}\n\nThis suite contains complex Word, Excel, PowerPoint, PDF, scanned PDF, image OCR, CSV, nested ZIP and error fixtures.\n\nImportant: The long-path test is generated at runtime so Windows can extract the project ZIP normally.\n\nTotal fixture files: {len(manifest)}\n'''
(OUT/'README.md').write_text(readme,encoding='utf-8')
print(f'Generated {len(manifest)} fixtures in {OUT}')
